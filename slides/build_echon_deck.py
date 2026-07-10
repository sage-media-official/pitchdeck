#!/usr/bin/env python3
"""Build an editable Echon Annual Marketing Plan deck (16:9) with python-pptx.
Everything is native text boxes / tables / autoshapes so it stays editable in
Google Slides after import."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import copy

# ---------- palette ----------
C = {
    'cream':'FDFBEE','cream2':'F6F3E2','cream3':'EFEBD6',
    'sage':'517354','sage6':'456245','sage7':'37503B','sage9':'22321F',
    'ink':'1C2620','mute':'646F5C',
    'echon':'C57B3C','echon2':'A85F2B','echonl':'E8B27A','gold':'D9A441',
    'white':'FFFFFF','line':'D9DBCB',
}
def rgb(h): return RGBColor.from_string(C.get(h,h))
TITLE='Glacial Indifference'
BODY='Inter'

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = prs.slide_width, prs.slide_height

def slide(bg='cream'):
    s = prs.slides.add_slide(BLANK)
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0,0, SW, SH)
    r.fill.solid(); r.fill.fore_color.rgb = rgb(bg); r.line.fill.background()
    r.shadow.inherit = False
    return s

def _set_font(run, size, color, font, bold, italic=False, spc=None):
    run.font.size = Pt(size); run.font.bold = bold; run.font.italic = italic
    run.font.name = font; run.font.color.rgb = rgb(color)
    rPr = run._r.get_or_add_rPr()
    if spc is not None:
        rPr.set('spc', str(int(spc)))  # character spacing in 1/100 pt (negative tightens)
    for tag in ('a:latin','a:cs'):
        e = rPr.find(qn(tag))
        if e is None:
            e = rPr.makeelement(qn(tag), {}); rPr.append(e)
        e.set('typeface', font)

def text(s, txt, l, t, w, h, size=14, color='ink', font=BODY, bold=False,
         align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, italic=False, sp_after=2,
         line_spacing=1.05, track=None):
    tb = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left=0; tf.margin_right=0; tf.margin_top=0; tf.margin_bottom=0
    lines = txt.split('\n') if isinstance(txt,str) else txt
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.alignment = align; p.space_after = Pt(sp_after); p.space_before = Pt(0)
        try: p.line_spacing = line_spacing
        except Exception: pass
        if isinstance(ln, list):  # list of (text,opts) runs
            for seg, opts in ln:
                r = p.add_run(); r.text = seg
                _set_font(r, opts.get('size',size), opts.get('color',color),
                          opts.get('font',font), opts.get('bold',bold), opts.get('italic',italic),
                          spc=opts.get('track',track))
        else:
            r = p.add_run(); r.text = ln
            _set_font(r, size, color, font, bold, italic, spc=track)
    return tb

def box(s, l, t, w, h, fill='white', line=None, line_w=1.0, radius=0.08, shadow=False, shape=MSO_SHAPE.ROUNDED_RECTANGLE):
    sp = s.shapes.add_shape(shape, Inches(l), Inches(t), Inches(w), Inches(h))
    if fill is None: sp.fill.background()
    else: sp.fill.solid(); sp.fill.fore_color.rgb = rgb(fill)
    if line is None: sp.line.fill.background()
    else: sp.line.color.rgb = rgb(line); sp.line.width = Pt(line_w)
    sp.shadow.inherit = False
    if shape==MSO_SHAPE.ROUNDED_RECTANGLE:
        try: sp.adjustments[0] = radius
        except Exception: pass
    if shadow:
        el = sp._element.spPr
        ef = el.makeelement(qn('a:effectLst'), {})
        sh = el.makeelement(qn('a:outerShdw'), {'blurRad':'120000','dist':'40000','dir':'5400000','rotWithShape':'0'})
        clr = el.makeelement(qn('a:srgbClr'), {'val':'22321F'})
        alpha = el.makeelement(qn('a:alpha'), {'val':'16000'})
        clr.append(alpha); sh.append(clr); ef.append(sh); el.append(ef)
    return sp

def shape_text(sp, txt, size=14, color='white', font=BODY, bold=False, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE, italic=False):
    tf = sp.text_frame; tf.word_wrap=True; tf.vertical_anchor=anchor
    tf.margin_left=Inches(0.12); tf.margin_right=Inches(0.12); tf.margin_top=Inches(0.04); tf.margin_bottom=Inches(0.04)
    lines = txt.split('\n')
    for i,ln in enumerate(lines):
        p = tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.alignment=align
        r = p.add_run(); r.text=ln
        _set_font(r, size, color, font, bold, italic)
    return sp

def eyebrow(s, n, label, l=0.7, t=0.6, dark=False):
    lbl = f"{n}   {label.upper()}"
    pw = 0.46 + 0.108*len(lbl)
    pill = box(s, l, t, pw, 0.4, fill=('1B2C24' if dark else 'F4E7D6'), line=None, radius=0.5)
    pill.fill.fore_color.rgb = rgb('2C3A2E') if dark else rgb('F6EBDC')
    tf = pill.text_frame; tf.word_wrap=False; tf.vertical_anchor=MSO_ANCHOR.MIDDLE
    tf.margin_left=Inches(0.16); tf.margin_right=Inches(0.16); tf.margin_top=Inches(0.02); tf.margin_bottom=Inches(0.02)
    p=tf.paragraphs[0]; p.alignment=PP_ALIGN.CENTER
    r=p.add_run(); r.text=lbl
    _set_font(r, 11, ('echonl' if dark else 'echon2'), BODY, True)
    return pill

def header(s, n, eyb, title, sub=None, dark=False, title_size=28, ty=1.18):
    # top accent bar
    bar = box(s, 0,0, SW, Inches(0.13), fill='echon', radius=0, shape=MSO_SHAPE.RECTANGLE)
    eyebrow(s, n, eyb, dark=dark)
    text(s, title, 0.7, ty, 12.0, 1.2, size=title_size, color=('white' if dark else 'sage9'), font=TITLE, bold=True, line_spacing=1.0, track=-30)
    if sub:
        text(s, sub, 0.7, ty+0.96, 11.6, 0.8, size=12.5, color=('cream3' if dark else 'mute'), font=BODY, line_spacing=1.18, track=-3)

def footer(s, n, dark=False):
    col = 'cream3' if dark else 'mute'
    text(s, "SÀGE MEDIA · Confidential — Echon (Kumar Arch Tech Ltd) Annual Marketing Plan",
         0.7, 7.06, 10.5, 0.35, size=8.5, color=col, font=BODY)
    text(s, str(n), 12.4, 7.06, 0.5, 0.35, size=8.5, color=col, font=BODY, align=PP_ALIGN.RIGHT)

# horizontal bar (rectangle) editable
def hbar(s, l, t, full_w, frac, h, fill, radius=0.4):
    box(s, l, t, full_w, h, fill='cream3', line=None, radius=radius)
    if frac>0:
        box(s, l, t, max(full_w*frac,0.06), h, fill=fill, line=None, radius=radius)

# vertical bar
def vbar(s, cx, base_y, w, height, fill, radius=0.18):
    box(s, cx-w/2, base_y-height, w, height, fill=fill, line=None, radius=radius, shape=MSO_SHAPE.ROUND_2_SAME_RECTANGLE if False else MSO_SHAPE.ROUNDED_RECTANGLE)

def style_table(tbl, header_fill='sage9', header_color='cream', body_size=10.5, header_size=10, zebra='cream2'):
    tblPr = tbl._tbl.tblPr
    # remove default style banding
    for el in list(tblPr):
        tblPr.remove(el)
    for ri, row in enumerate(tbl.rows):
        for ci, cell in enumerate(row.cells):
            cell.margin_left=Inches(0.1); cell.margin_right=Inches(0.1)
            cell.margin_top=Inches(0.04); cell.margin_bottom=Inches(0.04)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            if ri==0:
                cell.fill.solid(); cell.fill.fore_color.rgb = rgb(header_fill)
            else:
                cell.fill.solid(); cell.fill.fore_color.rgb = rgb(zebra if ri%2==0 else 'white')
    return tbl

def set_cell(cell, txt, size=10.5, color='ink', bold=False, align=PP_ALIGN.LEFT, font=BODY):
    tf = cell.text_frame; tf.word_wrap=True
    p = tf.paragraphs[0]; p.alignment=align
    for r in list(p.runs): r.text=''
    r = p.add_run(); r.text = txt
    _set_font(r, size, color, font, bold)

# =====================================================================
# SLIDE 1 — COVER
# =====================================================================
s = slide('cream')
box(s, 0,0, SW, Inches(0.13), fill='echon', radius=0, shape=MSO_SHAPE.RECTANGLE)
# right deco panel
box(s, 9.1, 0.13, 4.233, 7.37, fill='sage9', line=None, radius=0.0, shape=MSO_SHAPE.RECTANGLE)
box(s, 9.1, 0.13, 4.233, 7.37, fill='sage9', line=None, radius=0.0, shape=MSO_SHAPE.RECTANGLE)
# stacked "boards" motif on the dark panel
for i,(yy,cc) in enumerate([(1.2,'sage7'),(1.62,'sage6'),(2.04,'echon2'),(2.46,'echon'),(2.88,'gold')]):
    box(s, 9.7, yy, 3.1, 0.3, fill=cc, line=None, radius=0.3)
text(s, "PVC BLEND-BASED\nBUILDING MATERIALS", 9.7, 3.5, 3.2, 1.0, size=12, color='echonl', font=TITLE, bold=True, line_spacing=1.1)
text(s, "Boards · Doors · Wall &\nCeiling Panels · Flooring · Signage", 9.7, 4.35, 3.2, 1.0, size=11, color='cream3', font=BODY, line_spacing=1.2)
text(s, "Prepared by Sage Media\nJaipur · Mumbai · Dubai", 9.7, 6.2, 3.2, 0.8, size=10.5, color='cream2', font=BODY, line_spacing=1.2)

# left content
pill = box(s, 0.7, 0.7, 4.6, 0.44, fill='F6EBDC', line=None, radius=0.5)
shape_text(pill, "YEAR 1 · INDIA · ANNUAL MARKETING PLAN", size=11, color='echon2', bold=True, align=PP_ALIGN.CENTER)
text(s, "Echon", 0.66, 1.3, 8, 1.2, size=58, color='sage9', font=TITLE, bold=True, track=-50)
text(s, [
        [("Branding the worry-free ", {'color':'sage9'})],
        [("alternative to wood.", {'color':'echon2'})],
     ], 0.7, 2.5, 8.3, 1.6, size=30, color='sage9', font=TITLE, bold=True, line_spacing=1.05, track=-25)
text(s, "A sales-facilitation marketing engine that turns 22 years of manufacturing strength — and India's largest PVC-export base — into a branded domestic leader across channel, influencers & brand.",
     0.7, 4.05, 8.0, 1.2, size=13.5, color='mute', font=BODY, line_spacing=1.2)
# KPI mini cards
kpis = [("₹3.99 Cr","Year-1 investment"),("400","Dealers · 40 cities"),
        ("~8,000","Carpenters engaged"),("~₹70 Cr","Revenue enabled"),("17.5×","Return per ₹1")]
kx = 0.7; kw = 1.6; gap=0.04
for i,(v,l) in enumerate(kpis):
    cardw = (8.0-gap*4)/5
    cx = 0.7 + i*(cardw+gap)
    c = box(s, cx, 5.55, cardw, 1.15, fill='white', line='line', line_w=1, radius=0.12, shadow=True)
    text(s, v, cx, 5.7, cardw, 0.5, size=16.5, color='sage9', font=TITLE, bold=True, align=PP_ALIGN.CENTER)
    text(s, l, cx+0.05, 6.18, cardw-0.1, 0.5, size=8.5, color='mute', font=BODY, align=PP_ALIGN.CENTER, line_spacing=1.0)

# =====================================================================
# SLIDE 2 — EXECUTIVE SUMMARY
# =====================================================================
s = slide('sage9')
header(s, "00", "Executive Summary", "Convert manufacturing strength into branded demand.",
       "Deepen the dealer channel, win the influencer layer — carpenters above all — and create enough consumer pull that buyers ask for Echon by name. Every rupee is tied to a measurable outcome.",
       dark=True, title_size=30)
rows = [
    ("Recommended Investment","₹3.99 Cr / year (Balanced) · scenario range ₹1.85 – 7.90 Cr"),
    ("Year-1 Footprint","8 states · 40 cities · 400 dealers / distributors onboarded"),
    ("Influencer Base","~8,000 carpenters engaged · ~1,000 architects / interior designers"),
    ("Demand Generated","12,000 – 18,000 leads · ~2,000 qualified enquiries routed to sales"),
    ("Revenue Enabled","~₹70 Cr (~17.5× spend; marketing ≈ 5.7% of enabled revenue)"),
    ("Reporting","Monthly — against sales-facilitation KPIs, not vanity metrics"),
]
ry = 3.05
for k,v in rows:
    box(s, 0.7, ry+0.46, 8.0, 0.012, fill='456245', line=None, radius=0, shape=MSO_SHAPE.RECTANGLE)
    text(s, k.upper(), 0.7, ry, 2.7, 0.4, size=10.5, color='echonl', font=BODY, bold=True, anchor=MSO_ANCHOR.MIDDLE)
    text(s, v, 3.45, ry, 5.3, 0.45, size=11.5, color='cream', font=BODY, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.05)
    ry += 0.62
# highlight card
hc = box(s, 9.2, 2.95, 3.45, 3.5, fill='2C3A2E', line='456245', line_w=1, radius=0.08, shadow=False)
text(s, "BALANCED SCENARIO", 9.5, 3.2, 2.9, 0.4, size=11, color='echonl', font=BODY, bold=True)
text(s, "₹3.99 Cr", 9.5, 3.6, 2.9, 0.9, size=40, color='white', font=TITLE, bold=True)
text(s, "Our recommendation for a credible national Year-1 push.", 9.5, 4.6, 2.9, 0.7, size=11.5, color='cream3', font=BODY, line_spacing=1.15)
# range bar
box(s, 9.5, 5.5, 2.85, 0.16, fill='456245', line=None, radius=0.5)
box(s, 9.5+2.85*0.23, 5.5, 2.85*0.51, 0.16, fill='echon', line=None, radius=0.5)
text(s, "Conservative\n₹1.85 Cr", 9.5, 5.75, 1.4, 0.6, size=9.5, color='cream3', font=BODY, line_spacing=1.0)
text(s, "Aggressive\n₹7.90 Cr", 11.0, 5.75, 1.35, 0.6, size=9.5, color='cream3', font=BODY, align=PP_ALIGN.RIGHT, line_spacing=1.0)
footer(s, 1, dark=True)

# =====================================================================
# SLIDE 3 — STRATEGIC APPROACH (4 pillars)
# =====================================================================
s = slide('cream')
header(s, "01", "Strategic Approach", "The sale is decided in two places.",
       "A homeowner rarely names a board brand — the carpenter recommends it and the dealer stocks what moves. Four pillars, and every rupee maps to one.")
pillars = [
    ("01","Win the carpenter","The true decision-maker. Nukkad meets, a loyalty program, samples and recognition make Echon the default recommendation at the point of fabrication."),
    ("02","Arm the counter","Dealers & distributors are the revenue engine. Signage, counter branding, samples, meets and trade visibility help the sales team onboard and sell through."),
    ("03","Build consumer pull","ATL, always-on social, influencers, the visualizer and experience corners make consumers ask for Echon by name — easing the channel's job."),
    ("04","Convert credibility","Translate 22 years and 'India's largest exporter' into domestic trust — through PR, content, AID engagement and project specs, government & private."),
]
cw = 2.92; gap=0.18; x0=0.7; y0=3.0; ch=3.5
for i,(n,h,b) in enumerate(pillars):
    cx = x0 + i*(cw+gap)
    box(s, cx, y0, cw, ch, fill='white', line='line', line_w=1, radius=0.07, shadow=True)
    box(s, cx, y0+ch-0.1, cw, 0.1, fill='echon', line=None, radius=0, shape=MSO_SHAPE.RECTANGLE)
    text(s, f"PILLAR {n}", cx+0.25, y0+0.28, cw-0.5, 0.3, size=10, color='echon2', font=BODY, bold=True)
    ic = box(s, cx+0.25, y0+0.62, 0.66, 0.66, fill=('sage7' if i%2==0 else 'echon2'), line=None, radius=0.24)
    ic.text_frame.word_wrap=False
    shape_text(ic, str(i+1), size=24, color='cream', font=TITLE, bold=True, align=PP_ALIGN.CENTER)
    text(s, h, cx+0.25, y0+1.5, cw-0.5, 0.7, size=15.5, color='sage9', font=TITLE, bold=True, line_spacing=1.0)
    text(s, b, cx+0.25, y0+2.2, cw-0.5, 1.2, size=10.5, color='mute', font=BODY, line_spacing=1.12)
footer(s, 2)

# =====================================================================
# SLIDE 4 — MARKETING OBJECTIVES (table)
# =====================================================================
s = slide('cream')
header(s, "02", "Marketing Objectives", "Five objectives, five measurable Year-1 targets.")
objs = [
    ("1","Build brand salience","Establish Echon as the branded leader in PVC blend-based building materials — the new-age alternative to wood.","Recall & reach in 8 anchor markets"),
    ("2","Facilitate channel expansion","Help the sales team onboard and activate dealers / distributors and arm them to sell through.","400 dealers across 40 cities"),
    ("3","Win the influencer layer","Carpenters above all, plus architects, interior designers, contractors and builders.","8,000 carpenters · 1,000 AID"),
    ("4","Generate qualified demand","Dealer, retail, project (govt & private) and consumer enquiries that feed sales targets.","12,000–18,000 leads / year"),
    ("5","Run a measurable, ROI-driven engine","Reported monthly against sales facilitation, not vanity metrics.","~₹70 Cr revenue enabled"),
]
oy = 2.5; rh=0.83
for n,h,b,tgt in objs:
    box(s, 0.7, oy, 11.93, rh-0.12, fill='white', line='line', line_w=1, radius=0.1, shadow=True)
    text(s, n, 0.85, oy, 0.8, rh-0.12, size=30, color='echon', font=TITLE, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(s, h, 1.75, oy+0.1, 6.4, 0.4, size=14, color='sage9', font=TITLE, bold=True)
    text(s, b, 1.75, oy+0.42, 6.7, 0.4, size=10.5, color='mute', font=BODY, line_spacing=1.05)
    box(s, 8.7, oy+0.13, 0.02, rh-0.38, fill='line', line=None, radius=0, shape=MSO_SHAPE.RECTANGLE)
    text(s, tgt, 8.95, oy, 3.5, rh-0.12, size=12.5, color='sage7', font=BODY, bold=True, anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.LEFT, line_spacing=1.0)
    oy += rh
footer(s, 3)

# =====================================================================
# SLIDE 5 — TARGET SEGMENTS
# =====================================================================
s = slide('cream2')
header(s, "03", "Target Segments", "Three segments, three jobs.",
       "The channel sells, the influencers recommend, the consumer pulls. Marketing serves all three at once.")
segs = [
    ("SEGMENT 1","Channel & Direct Sales","The revenue engine","sage7",
     "Dealers & distributors, retailers and large projects — government and private. The stock points that move volume.",
     "Trade visibility, dealer/contractor meets, loyalty, samples & catalogs, exhibition leads, project-spec support, corporate AV.",
     "Onboard, stock, sell-through"),
    ("SEGMENT 2","The Influencers","The recommenders","echon2",
     "Architects, interior designers, builders/contractors — and most importantly carpenters, who specify and fabricate.",
     "Carpenter nukkad meets + loyalty, AID engagement & CPD, visualizer / 3D plugin, technical content, samples & recognition.",
     "Become the default recommendation"),
    ("SEGMENT 3","End Consumers","The pull","gold",
     "Homeowners, renovators and self-builders choosing durable, low-maintenance surfaces over wood.",
     "ATL, always-on social, influencer content, website + visualizer, performance lead-gen, experience corners, brand ambassador.",
     "Create pull"),
]
cw=3.9; gap=0.2; x0=0.7; y0=3.05; ch=3.6
for i,(tag,h,role,col,who,win,goal) in enumerate(segs):
    cx = x0 + i*(cw+gap)
    box(s, cx, y0, cw, ch, fill='white', line='line', line_w=1, radius=0.07, shadow=True)
    tp = box(s, cx+0.25, y0-0.18, 1.5, 0.38, fill=col, line=None, radius=0.5)
    shape_text(tp, tag, size=9.5, color='white', bold=True, align=PP_ALIGN.CENTER)
    text(s, h, cx+0.25, y0+0.32, cw-0.5, 0.5, size=15.5, color='sage9', font=TITLE, bold=True)
    text(s, role.upper(), cx+0.25, y0+0.78, cw-0.5, 0.3, size=9.5, color='echon2', font=BODY, bold=True)
    text(s, who, cx+0.25, y0+1.12, cw-0.5, 0.9, size=10.5, color='ink', font=BODY, line_spacing=1.1)
    box(s, cx+0.25, y0+2.15, cw-0.5, 0.012, fill='line', line=None, radius=0, shape=MSO_SHAPE.RECTANGLE)
    text(s, win, cx+0.25, y0+2.25, cw-0.5, 0.9, size=9.8, color='mute', font=BODY, line_spacing=1.08)
    gp = box(s, cx+0.25, y0+ch-0.55, cw-0.5, 0.4, fill='cream3', line=None, radius=0.5)
    shape_text(gp, "Goal: "+goal, size=10, color='sage7', bold=True, align=PP_ALIGN.CENTER)
footer(s, 4)

# =====================================================================
# SLIDE 6 — SALES ROLLOUT (table + dealer bars)
# =====================================================================
s = slide('cream')
header(s, "04", "Sales Rollout", "A phased Year-1 footprint, front-loaded into the anchor states.",
       "8 states · 40 priority cities · 400 dealers / distributors. Proposed — to be confirmed with Echon's sales leadership.")
data = [
    ("Rajasthan (home)","Anchor","6","60","Q1"),
    ("Gujarat","Anchor","6","60","Q1"),
    ("Maharashtra","Anchor","7","70","Q1"),
    ("Delhi-NCR","Anchor","5","55","Q2"),
    ("Uttar Pradesh","Growth","6","50","Q2"),
    ("Punjab & Haryana","Growth","4","35","Q2"),
    ("Karnataka","Growth","3","35","Q3"),
    ("Telangana","Growth","3","35","Q3"),
    ("TOTAL — 8 states","Q1–Q3","40","400","—"),
]
rows_n=len(data)+1
tbl_shape = s.shapes.add_table(rows_n, 5, Inches(0.7), Inches(3.0), Inches(6.5), Inches(3.6))
tbl = tbl_shape.table
tbl.columns[0].width=Inches(2.5); tbl.columns[1].width=Inches(1.1); tbl.columns[2].width=Inches(0.9); tbl.columns[3].width=Inches(1.1); tbl.columns[4].width=Inches(0.9)
hdr=["State","Tier","Cities","Dealers","Activ."]
for ci,htxt in enumerate(hdr):
    set_cell(tbl.cell(0,ci), htxt, size=10, color='cream', bold=True, align=(PP_ALIGN.LEFT if ci==0 else PP_ALIGN.CENTER))
for ri,row in enumerate(data, start=1):
    last = ri==len(data)
    for ci,val in enumerate(row):
        set_cell(tbl.cell(ri,ci), val, size=10, color=('white' if last else 'ink'),
                 bold=last or ci==3, align=(PP_ALIGN.LEFT if ci==0 else PP_ALIGN.CENTER))
style_table(tbl)
# color rows
for ri in range(rows_n):
    for ci in range(5):
        cell=tbl.cell(ri,ci)
        if ri==0: cell.fill.solid(); cell.fill.fore_color.rgb=rgb('sage9')
        elif ri==len(data): cell.fill.solid(); cell.fill.fore_color.rgb=rgb('sage9')
        else: cell.fill.solid(); cell.fill.fore_color.rgb=rgb('cream2' if ri%2==0 else 'white')
text(s, "Q4 is consolidation & depth — no new states. Deepen the 8 active states, expand retail width, drive repeat orders, and pilot East India (West Bengal) ahead of Year 2.",
     0.7, 6.65, 6.5, 0.5, size=9.5, color='mute', font=BODY, italic=True, line_spacing=1.05)
# dealer bar chart (right)
box(s, 7.55, 3.0, 5.08, 3.85, fill='white', line='line', line_w=1, radius=0.06, shadow=True)
text(s, "Target dealers by state", 7.8, 3.2, 4.6, 0.4, size=13.5, color='sage9', font=TITLE, bold=True)
text(s, "Anchor states (green) front-loaded; growth states (copper) follow.", 7.8, 3.58, 4.6, 0.4, size=9.5, color='mute', font=BODY)
bars=[("Maharashtra",70,'sage7'),("Rajasthan",60,'sage7'),("Gujarat",60,'sage7'),("Delhi-NCR",55,'sage7'),
      ("Uttar Pradesh",50,'echon'),("Punjab & Hy.",35,'echon'),("Karnataka",35,'echon'),("Telangana",35,'echon')]
by=4.1; bh=0.24; bgap=0.105; bx=9.05; bfull=2.9
for name,val,col in bars:
    text(s, name, 7.78, by-0.02, 1.25, 0.3, size=9, color='ink', font=BODY, align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
    hbar(s, bx, by, bfull, val/70.0, bh, col)
    text(s, str(val), bx+bfull+0.08, by-0.02, 0.45, 0.3, size=9.5, color='sage9', font=BODY, bold=True, anchor=MSO_ANCHOR.MIDDLE)
    by += bh+bgap
footer(s, 5)

# =====================================================================
# SLIDE 7 — PLAN OF ACTION (quarterly spend chart + quarters)
# =====================================================================
s = slide('cream2')
header(s, "05", "Plan of Action — 12 Months", "Four quarters: build, expand, deepen, optimise.",
       "Spend is paced to match the rollout — front-loaded into the build, then sustained.")
# quarterly bar chart
box(s, 0.7, 3.0, 6.2, 3.85, fill='white', line='line', line_w=1, radius=0.06, shadow=True)
text(s, "Quarterly spend phasing — Balanced (₹ Lakhs)", 0.95, 3.2, 5.7, 0.4, size=13, color='sage9', font=TITLE, bold=True)
qd=[("Q1",140,'echon'),("Q2",110,'sage7'),("Q3",88,'sage7'),("Q4",61,'sage7')]
base_y=6.1; maxh=2.0; cw2=1.05
xs=[1.45,2.85,4.25,5.65]
for (q,v,col),cx in zip(qd,xs):
    hgt=maxh*v/140.0
    box(s, cx-cw2/2, base_y-hgt, cw2, hgt, fill=col, line=None, radius=0.12)
    text(s, f"₹{v} L", cx-0.6, base_y-hgt-0.32, 1.2, 0.3, size=11, color='sage9', font=BODY, bold=True, align=PP_ALIGN.CENTER)
    text(s, q, cx-0.6, base_y+0.05, 1.2, 0.3, size=12, color='sage9', font=TITLE, bold=True, align=PP_ALIGN.CENTER)
text(s, "Foundation        N/W expansion        South & depth        Optimise", 0.95, 6.5, 5.7, 0.3, size=8.5, color='mute', font=BODY)
# quarter cards (right)
qcards=[("Q1 · M1–M3","Foundation & Anchor Launch","Set up anchor teams (RJ/GJ/MH); brand toolkit, website + visualizer, loyalty portal; signage & first carpenter nukkad meets. ~120 dealers active."),
        ("Q2 · M4–M6","North / West Expansion","Activate Delhi-NCR, UP, Punjab/Haryana (~230 dealers). ATL burst, AID program, first regional exhibition, mid-year sell-through review."),
        ("Q3 · M7–M9","South Expansion & Depth","Activate Karnataka & Telangana (~330 dealers). Project & AID specs, second expo, experience corner, performance optimisation."),
        ("Q4 · M10–M12","Consolidate & Optimise","Retail width within ~400 dealers, repeat-order & festive push, WB pilot prep, brand-ambassador decision, Year-2 roadmap.")]
qy=3.0; qh=0.93
for tag,h,b in qcards:
    box(s, 7.25, qy, 5.4, qh-0.1, fill='white', line='line', line_w=1, radius=0.09, shadow=True)
    box(s, 7.25, qy, 0.11, qh-0.1, fill='echon', line=None, radius=0, shape=MSO_SHAPE.RECTANGLE)
    text(s, tag, 7.55, qy+0.1, 2.3, 0.3, size=10, color='echon2', font=BODY, bold=True)
    text(s, h, 9.0, qy+0.08, 3.5, 0.3, size=12, color='sage9', font=TITLE, bold=True)
    text(s, b, 7.55, qy+0.4, 4.9, 0.5, size=9, color='mute', font=BODY, line_spacing=1.04)
    qy+=qh
footer(s, 6)

# =====================================================================
# SLIDE 8 — 12-MONTH PLAN TABLE
# =====================================================================
s = slide('cream')
header(s, "05", "Plan of Action · Calendar", "The Year-1 rollout calendar.")
months=[
 ("M1","Set up anchor teams; map RJ/GJ/MH dealers","Brand toolkit; website + visualizer & 3D plugin; social revamp; WhatsApp + loyalty setup"),
 ("M2","Onboard first anchor dealers (RJ/GJ/MH)","Signage + counter branding; first carpenter meets; SEO + blogs; performance pilot; samples"),
 ("M3","~120 anchor dealers active","Dealer/contractor meets; influencer collabs start; product films + hero DVC; expo prep"),
 ("M4","Activate Delhi-NCR; add dealers","ATL burst (print + OOH); scale performance; AID program; visualizer-led spec push"),
 ("M5","Activate UP + Punjab/Haryana (~230)","Signage/counter in new cities; carpenter meets; first regional exhibition; sample wave 2"),
 ("M6","Mid-year sell-through review","Regional TVC/radio burst; loyalty scale-up; mid-year performance review & creative refresh"),
 ("M7","Activate Karnataka; project & AID specs","Influencer push; experience corner setup; walk-through content; retargeting layer"),
 ("M8","Activate Telangana (~330)","Second expo (Acetech-class); visualizer-led AID campaign; CPL reduction focus"),
 ("M9","Depth in active states","Counter-branding depth; carpenter loyalty milestones; case studies, corporate AV"),
 ("M10","Retail width within dealers (~400)","Repeat-order & festive campaign; merchandise refresh; project-spec content; remarketing"),
 ("M11","East India (WB) pilot prep","Full-funnel optimisation; loyalty redemption drive; PR/awards; influencer round 2"),
 ("M12","Annual sales vs target review","Year-in-review content; brand-ambassador decision; KPI review; Year-2 roadmap & budget"),
]
tbl_shape=s.shapes.add_table(13,3,Inches(0.7),Inches(2.05),Inches(11.93),Inches(4.9))
tbl=tbl_shape.table
tbl.columns[0].width=Inches(0.85); tbl.columns[1].width=Inches(4.0); tbl.columns[2].width=Inches(7.08)
for ci,htxt in enumerate(["Mo.","Sales / Rollout Focus","Key Marketing Actions"]):
    set_cell(tbl.cell(0,ci),htxt,size=10.5,color='cream',bold=True)
for ri,(m,a,b) in enumerate(months,start=1):
    set_cell(tbl.cell(ri,0),m,size=10,color='echon2',bold=True,align=PP_ALIGN.CENTER)
    set_cell(tbl.cell(ri,1),a,size=9.5,color='sage9',bold=True)
    set_cell(tbl.cell(ri,2),b,size=9.5,color='mute')
for ri in range(13):
    for ci in range(3):
        cell=tbl.cell(ri,ci)
        cell.fill.solid()
        cell.fill.fore_color.rgb=rgb('sage9') if ri==0 else rgb('cream2' if ri%2==0 else 'white')
        cell.margin_top=Inches(0.02); cell.margin_bottom=Inches(0.02)
footer(s, 7)

# =====================================================================
# SLIDE 9 — CHANNEL STRATEGY & KPIs (table)
# =====================================================================
s = slide('cream')
header(s, "06", "Channel Strategy & KPIs", "How each channel works — and what it returns.",
       "Balanced working-media spend (₹ L, ex-agency); a flat 20% agency fee applies on the budget.")
ch=[
 ("ATL — Print/OOH/Radio/TVC","55","~25–40M reach across anchor cities; trade & architect press; dealer-hub OOH"),
 ("Digital — Social/Influencer/Web/SEO","68","+40–60k followers; 12–18M reach; 60–80 collabs; 30–50k visualizer sessions"),
 ("Performance — SEM/SMM/GDN","38","12,000–18,000 leads/yr; CPL ₹120–250; ~1,500–2,500 qualified enquiries"),
 ("Content / Videos","30","120–180 assets; 8–12M video views; product/application film + walk-throughs"),
 ("Exhibition","17","4–6 expos; 2,000–3,500 leads; 60–100 dealer/distributor signups"),
 ("Visibility — Signage + Counter","23","~400 outlets with signage; 250–350 counter-branded points across 40 cities"),
 ("Meets — Carpenter + Dealer","28","~320 carpenter nukkad meets → ~8,000 carpenters; ~80 dealer/contractor meets"),
 ("Loyalty — Dealer & Carpenter","22","8,000+ carpenters & 400+ dealers enrolled; 15–25% repeat-purchase uplift"),
 ("Samples + Collaterals","22","~12,000–15,000 sample kits; full sales-kit & sample-book coverage"),
 ("Experience Stores / Centers","16","1–2 experience corners in anchor cities; touch-&-feel for AID & consumers"),
 ("Agency Fee — Mainline + Digital","20%","20% of budget (₹80 L Bal): strategy, creative, execution & monthly reporting"),
]
tbl_shape=s.shapes.add_table(len(ch)+1,3,Inches(0.7),Inches(2.5),Inches(11.93),Inches(4.4))
tbl=tbl_shape.table
tbl.columns[0].width=Inches(3.7); tbl.columns[1].width=Inches(1.2); tbl.columns[2].width=Inches(7.03)
for ci,htxt in enumerate(["Channel","₹ L","Primary KPIs & Expected Year-1 Numbers"]):
    set_cell(tbl.cell(0,ci),htxt,size=10.5,color='cream',bold=True,align=(PP_ALIGN.CENTER if ci==1 else PP_ALIGN.LEFT))
for ri,(c,bud,k) in enumerate(ch,start=1):
    set_cell(tbl.cell(ri,0),c,size=10,color='sage9',bold=True)
    set_cell(tbl.cell(ri,1),bud,size=10.5,color='echon2',bold=True,align=PP_ALIGN.CENTER)
    set_cell(tbl.cell(ri,2),k,size=9.5,color='mute')
for ri in range(len(ch)+1):
    for ci in range(3):
        cell=tbl.cell(ri,ci); cell.fill.solid()
        cell.fill.fore_color.rgb=rgb('sage9') if ri==0 else rgb('cream2' if ri%2==0 else 'white')
        cell.margin_top=Inches(0.015); cell.margin_bottom=Inches(0.015)
footer(s, 8)

# =====================================================================
# SLIDE 10 — BUDGET ALLOCATION (scenarios + stacked allocation)
# =====================================================================
s = slide('cream2')
header(s, "07", "Budget Allocation", "Match spend to ambition.",
       "Three scenarios. Balanced (₹3.99 Cr) is our recommendation — deliberately channel- and influencer-heavy.")
sc=[("CONSERVATIVE","₹1.85 Cr","A focused, anchor-first build. Proves the model before scaling.",'white','line','sage9'),
    ("BALANCED  ★","₹3.99 Cr","A credible national Year-1 push — all engines funded. Recommended.",'white','echon','sage9'),
    ("AGGRESSIVE","₹7.90 Cr","Adds a brand ambassador, full TVC burst & richer experience centres.",'white','line','sage9')]
cw=3.9; gap=0.2; x0=0.7; y0=3.0
for i,(nm,val,desc,fill,ln,vc) in enumerate(sc):
    cx=x0+i*(cw+gap)
    rec = i==1
    box(s, cx, y0, cw, 1.85, fill=fill, line=ln, line_w=(2 if rec else 1), radius=0.08, shadow=rec)
    text(s, nm, cx+0.3, y0+0.22, cw-0.6, 0.3, size=11, color=('echon2' if rec else 'mute'), font=BODY, bold=True)
    text(s, val, cx+0.3, y0+0.55, cw-0.6, 0.6, size=30, color=vc, font=TITLE, bold=True)
    text(s, desc, cx+0.3, y0+1.25, cw-0.6, 0.6, size=10, color='mute', font=BODY, line_spacing=1.08)
text(s, "Each total is all-in: 80% working media + 20% agency fee.  Balanced = ₹3.19 Cr media + ₹0.80 Cr agency = ₹3.99 Cr.",
     0.7, 4.95, 11.9, 0.3, size=10, color='echon2', font=BODY, bold=True)
# Where the money goes — 100% stacked bar (editable rectangles)
text(s, "Where the money goes — Balanced", 0.7, 5.2, 7, 0.4, size=14, color='sage9', font=TITLE, bold=True)
segs=[("Channel, Trade & Influencer",32,'sage7'),("Demand Generation",27,'echon'),
      ("Agency fee (20%)",20,'mute'),("ATL / Brand",14,'gold'),("Content",7,'sage')]
bx=0.7; bw=11.93; by=5.7; bh=0.6; acc=0
for nm,pct,col in segs:
    w=bw*pct/100.0
    seg=box(s, bx+acc, by, w, bh, fill=col, line='cream2', line_w=1.5, radius=0.0, shape=MSO_SHAPE.RECTANGLE)
    shape_text(seg, f"{pct:.0f}%", size=12, color=('sage9' if col=='gold' else 'white'), bold=True, align=PP_ALIGN.CENTER)
    acc+=w
# legend
lx=0.7; ly=6.5
for nm,pct,col in segs:
    box(s, lx, ly+0.05, 0.18,0.18, fill=col, line='line', line_w=0.75, radius=0.2)
    t=text(s, nm, lx+0.26, ly, 2.4, 0.4, size=9, color='ink', font=BODY, line_spacing=1.0)
    lx += 0.26 + (0.07*len(nm)) + 0.35
footer(s, 9)

# =====================================================================
# SLIDE 11 — BUDGET BY CHANNEL (bars + table)
# =====================================================================
s = slide('cream')
header(s, "07", "Budget — By Channel & Scenario", "Working-media allocation (80%); agency fee is a flat 20% of budget.")
# left: horizontal bars
box(s, 0.7, 2.5, 5.6, 4.35, fill='white', line='line', line_w=1, radius=0.06, shadow=True)
text(s, "Working media by channel — Balanced (₹ L)", 0.95, 2.68, 5.2, 0.4, size=12, color='sage9', font=TITLE, bold=True)
bb=[("Digital",68,'sage7'),("ATL",55,'sage7'),("Performance",38,'sage7'),("Content",30,'sage7'),
    ("Meets",28,'echon'),("Visibility",23,'echon'),("Loyalty",22,'echon'),("Exhibition",17,'echon'),
    ("Experience",16,'echon'),("Samples",13,'echon'),("Collaterals",9,'echon')]
by=3.2; bh=0.22; bgap=0.088; bx=2.0; bfull=3.4
for nm,v,col in bb:
    text(s, nm, 0.85, by-0.03, 1.05, 0.3, size=8.5, color='ink', font=BODY, align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
    hbar(s, bx, by, bfull, v/68.0, bh, col)
    text(s, str(v), bx+bfull+0.06, by-0.03, 0.4, 0.3, size=8.5, color='sage9', font=BODY, bold=True, anchor=MSO_ANCHOR.MIDDLE)
    by+=bh+bgap
# right: scenario table (media categories + agency + totals)
catf=[("ATL",120,55,12),("Digital",111,68,39),("Performance Mktg",64,38,19),("Content / Videos",54,30,13),
      ("Exhibition",30,17,9),("Visibility",38,23,13),("Meets",43,28,18),("Loyalty",34,22,12),
      ("Collaterals/Merch",15,9,5),("Samples",21,13,8),("Experience",34,16,0),("Brand Ambassador",68,0,0)]
specials=[("Working media (₹L)",632,319,148),("Agency fee — 20%",158,80,37),("Total (₹ Cr)","7.90","3.99","1.85")]
nrows=1+len(catf)+len(specials)
tbl_shape=s.shapes.add_table(nrows,4,Inches(6.55),Inches(2.5),Inches(6.08),Inches(4.35))
tbl=tbl_shape.table
tbl.columns[0].width=Inches(2.5); tbl.columns[1].width=Inches(1.19); tbl.columns[2].width=Inches(1.2); tbl.columns[3].width=Inches(1.19)
for ci,htxt in enumerate(["Category (media)","Aggr.","Bal.","Cons."]):
    set_cell(tbl.cell(0,ci),htxt,size=9,color='cream',bold=True,align=(PP_ALIGN.LEFT if ci==0 else PP_ALIGN.RIGHT))
def fnum(v): return "–" if v==0 else str(v)
for ri,(c,a,b,co) in enumerate(catf,start=1):
    set_cell(tbl.cell(ri,0),c,size=8.5,color='ink')
    set_cell(tbl.cell(ri,1),fnum(a),size=8.5,color='ink',align=PP_ALIGN.RIGHT)
    set_cell(tbl.cell(ri,2),fnum(b),size=8.5,color='sage9',bold=True,align=PP_ALIGN.RIGHT)
    set_cell(tbl.cell(ri,3),fnum(co),size=8.5,color='mute',align=PP_ALIGN.RIGHT)
mediarow=1+len(catf); agencyrow=mediarow+1; totalrow=agencyrow+1
for off,(c,a,b,co) in enumerate(specials):
    ri=mediarow+off
    dark = c.startswith('Working') or c.startswith('Total')
    set_cell(tbl.cell(ri,0),c,size=9,color=('cream' if dark else 'echon2'),bold=True)
    set_cell(tbl.cell(ri,1),str(a),size=9,color=('cream' if dark else 'echon2'),bold=True,align=PP_ALIGN.RIGHT)
    set_cell(tbl.cell(ri,2),str(b),size=9,color=('echonl' if dark else 'echon2'),bold=True,align=PP_ALIGN.RIGHT)
    set_cell(tbl.cell(ri,3),str(co),size=9,color=('cream' if dark else 'echon2'),bold=True,align=PP_ALIGN.RIGHT)
for ri in range(nrows):
    for ci in range(4):
        cell=tbl.cell(ri,ci); cell.fill.solid()
        if ri==0 or ri==mediarow or ri==totalrow: cell.fill.fore_color.rgb=rgb('sage9')
        elif ri==agencyrow: cell.fill.fore_color.rgb=rgb('cream3')
        else: cell.fill.fore_color.rgb=rgb('cream2' if ri%2==0 else 'white')
        cell.margin_top=Inches(0.006); cell.margin_bottom=Inches(0.006)
footer(s, 10)

# =====================================================================
# SLIDE 12 — THREE SCENARIOS COMPARED (grouped columns)
# =====================================================================
s = slide('cream2')
header(s, "07", "The Three Scenarios Compared", "By strategic bucket (₹ Lakhs).",
       "Channel & influencer spend scales hardest from Conservative → Aggressive; the agency fee stays a flat 20%.")
box(s, 0.7, 3.0, 11.93, 3.85, fill='white', line='line', line_w=1, radius=0.05, shadow=True)
groups=[("Channel &\nInfluencer",283,128,65),("Demand Gen",175,106,58),("Agency\n(20%)",158,80,37),("ATL / Brand",120,55,12),("Content",54,30,13)]
gx0=1.7; gspan=2.25; base_y=6.0; maxh=2.4; maxv=283; bw=0.5
for gi,(nm,a,b,co) in enumerate(groups):
    gx=gx0+gi*gspan
    for j,(v,col) in enumerate([(a,'echon2'),(b,'sage7'),(co,'cream3')]):
        cx=gx+j*(bw+0.06)
        hgt=maxh*v/maxv if v>0 else 0.03
        bxs=box(s, cx, base_y-hgt, bw, max(hgt,0.03), fill=col, line=('line' if col=='cream3' else None), line_w=0.75, radius=0.12)
        text(s, str(v) if v>0 else "–", cx-0.15, base_y-hgt-0.26, bw+0.3, 0.25, size=8, color='sage9', font=BODY, bold=True, align=PP_ALIGN.CENTER)
    text(s, nm, gx-0.2, base_y+0.08, bw*3+0.2, 0.5, size=10, color='sage9', font=BODY, bold=True, align=PP_ALIGN.CENTER, line_spacing=0.95)
# legend
lz=[("Aggressive · ₹7.90 Cr",'echon2'),("Balanced · ₹3.99 Cr",'sage7'),("Conservative · ₹1.85 Cr",'cream3')]
lx=3.0
for nm,col in lz:
    box(s, lx, 6.55, 0.2,0.2, fill=col, line='line', line_w=0.75, radius=0.2)
    text(s, nm, lx+0.28, 6.52, 2.6, 0.3, size=9.5, color='mute', font=BODY)
    lx+=2.9
footer(s, 11)

# =====================================================================
# SLIDE 13 — RETURN: DEMAND FUNNEL + COST PER OUTCOME
# =====================================================================
s = slide('cream')
header(s, "08", "Return on the Spend", "Demand created and revenue enabled — not impressions.",
       "A transparent, conservative model. Every input can be adjusted with Echon's actual numbers.")
# funnel (trapezoids)
text(s, "The Year-1 demand funnel", 0.7, 2.95, 6, 0.4, size=13.5, color='sage9', font=TITLE, bold=True)
fstages=[("12,000–18,000","Leads generated · paid + digital",'sage7',6.2),
         ("~2,000","Qualified enquiries routed to sales",'sage',5.0),
         ("~250","Project / retail orders · ~₹4 L avg",'echon',3.8),
         ("~₹70 Cr","Revenue enabled · channel + direct",'echon2',2.6)]
fy=3.45; fx_c=3.55
for val,lab,col,w in fstages:
    tp=box(s, fx_c-w/2, fy, w, 0.62, fill=col, line=None, radius=0, shape=MSO_SHAPE.TRAPEZOID)
    tp.rotation=180
    # text overlaid (separate textbox so it's upright)
    text(s, val, fx_c-w/2, fy+0.06, w, 0.3, size=14, color='white', font=BODY, bold=True, align=PP_ALIGN.CENTER)
    text(s, lab, fx_c-w/2, fy+0.34, w, 0.25, size=8.5, color='cream', font=BODY, align=PP_ALIGN.CENTER)
    fy+=0.72
# revenue side notes
sidenotes=[("₹60 Cr","400 dealers @ ~₹15 L ramped offtake"),("₹10 Cr","~250 project/retail orders @ ~₹4 L"),("~₹21 Cr","Gross profit at ~30% margin")]
sy=3.5
for v,l in sidenotes:
    box(s, 7.4, sy, 5.2, 0.78, fill='cream2', line='line', line_w=1, radius=0.1)
    text(s, v, 7.6, sy+0.13, 1.8, 0.5, size=20, color='echon2', font=TITLE, bold=True, anchor=MSO_ANCHOR.MIDDLE)
    text(s, l, 9.3, sy+0.13, 3.1, 0.5, size=10, color='mute', font=BODY, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0)
    sy+=0.9
# cost per outcome strip
text(s, "Cost per outcome (Balanced)", 0.7, 6.05, 6, 0.3, size=11.5, color='sage9', font=TITLE, bold=True)
cpo=[("₹150–250","per lead"),("~₹2,000","per qualified enquiry"),("~₹32,000","per dealer armed"),("~₹375","per carpenter"),("~5.7%","of revenue enabled")]
cx=0.7
for v,l in cpo:
    box(s, cx, 6.4, 2.35, 0.55, fill='white', line='line', line_w=1, radius=0.12, shadow=True)
    text(s, v, cx+0.12, 6.46, 1.2, 0.45, size=13, color='sage9', font=TITLE, bold=True, anchor=MSO_ANCHOR.MIDDLE)
    text(s, l, cx+1.15, 6.46, 1.15, 0.45, size=8, color='mute', font=BODY, anchor=MSO_ANCHOR.MIDDLE, line_spacing=0.95)
    cx+=2.42
footer(s, 12)

# =====================================================================
# SLIDE 14 — REVENUE VS SPEND (ROI)
# =====================================================================
s = slide('sage9')
header(s, "08", "Revenue Enabled vs. Spend", "~₹70 Cr enabled against ₹3.99 Cr spend.",
       "Two engines: 400 new dealers (≈ ₹60 Cr) and ~250 direct project/retail orders (≈ ₹10 Cr).", dark=True)
# bars
roidata=[("Spend","₹3.99 Cr",0.057,'cream3'),("Gross profit","₹21 Cr",0.30,'sage'),("Revenue enabled","₹70 Cr",1.0,'echon')]
base_y=6.0; maxh=2.6; xs=[1.6,3.3,5.0]; bw=1.1
for (nm,lab,frac,col),cx in zip(roidata,xs):
    hgt=maxh*frac
    box(s, cx-bw/2, base_y-hgt, bw, hgt, fill=col, line=None, radius=0.1)
    text(s, lab, cx-0.8, base_y-hgt-0.34, 1.6, 0.3, size=13, color='white', font=BODY, bold=True, align=PP_ALIGN.CENTER)
    text(s, nm, cx-0.8, base_y+0.08, 1.6, 0.3, size=10.5, color='cream3', font=BODY, align=PP_ALIGN.CENTER)
# big 17.5x
box(s, 7.3, 3.1, 5.35, 3.4, fill='2C3A2E', line='456245', line_w=1, radius=0.08)
text(s, "17.5×", 7.5, 3.35, 5.0, 1.4, size=72, color='echon', font=TITLE, bold=True, align=PP_ALIGN.CENTER)
text(s, "revenue enabled per ₹1 spent", 7.5, 4.85, 5.0, 0.4, size=14, color='cream', font=BODY, align=PP_ALIGN.CENTER)
p1=box(s, 7.9, 5.45, 4.15, 0.42, fill='456245', line=None, radius=0.5)
shape_text(p1, "~5.7% of revenue (build-year)", size=10.5, color='cream', bold=True, align=PP_ALIGN.CENTER)
p2=box(s, 7.9, 5.95, 4.15, 0.42, fill='456245', line=None, radius=0.5)
shape_text(p2, "Matures to 2–3% as the channel compounds", size=10.5, color='cream', bold=True, align=PP_ALIGN.CENTER)
footer(s, 13, dark=True)

# =====================================================================
# SLIDE 15 — GOVERNANCE + NEXT STEPS
# =====================================================================
s = slide('cream')
header(s, "09", "Governance & Next Steps", "Reported monthly against sales facilitation — not vanity metrics.")
gov=[("Demand","Leads, CPL, qualified dealer / project / consumer enquiries routed to sales."),
     ("Channel","Dealers onboarded vs target, outlets branded, sell-through, repeat orders."),
     ("Influencers","Carpenters & AID engaged / enrolled, loyalty active rate, specifications."),
     ("Brand","Reach, followers, engagement, recall in anchor markets."),
     ("Efficiency","Cost per lead, per dealer, per carpenter; revenue enabled vs spend.")]
gw=2.3; gap=0.12; x0=0.7; y0=2.45
for i,(h,b) in enumerate(gov):
    cx=x0+i*(gw+gap)
    box(s, cx, y0, gw, 1.7, fill='white', line='line', line_w=1, radius=0.08, shadow=True)
    box(s, cx, y0, gw, 0.09, fill='echon', line=None, radius=0, shape=MSO_SHAPE.RECTANGLE)
    text(s, h, cx+0.2, y0+0.25, gw-0.4, 0.35, size=13, color='sage9', font=TITLE, bold=True)
    text(s, b, cx+0.2, y0+0.65, gw-0.4, 1.0, size=9.5, color='mute', font=BODY, line_spacing=1.1)
# next steps
text(s, "Next Steps", 0.7, 4.5, 6, 0.4, size=16, color='sage9', font=TITLE, bold=True)
steps=[("01","Lock the plan","Confirm the budget scenario and lock the state → city → dealer rollout with Echon's sales leadership."),
       ("02","Kick-off (M1)","Brand toolkit, website + visualizer, loyalty portal, and the channel onboarding kit."),
       ("03","Go live (M2)","Activate anchor states (RJ/GJ/MH) — first leads, meets and dealer onboarding in motion.")]
sw=3.9; gap=0.2; x0=0.7; y0=5.05
for n,h,b in steps:
    cx=x0+ (steps.index((n,h,b)))*(sw+gap)
    box(s, cx, y0, sw, 1.55, fill='cream2', line='line', line_w=1, radius=0.08, shadow=True)
    text(s, n, cx+0.25, y0+0.18, 1.0, 0.5, size=26, color='echonl', font=TITLE, bold=True)
    text(s, h, cx+0.25, y0+0.72, sw-0.5, 0.3, size=13, color='sage9', font=TITLE, bold=True)
    text(s, b, cx+0.25, y0+1.05, sw-0.5, 0.5, size=9.5, color='mute', font=BODY, line_spacing=1.05)
footer(s, 14)

# =====================================================================
# SLIDE 16 — CLOSING
# =====================================================================
s = slide('sage9')
box(s, 0,0, SW, Inches(0.13), fill='echon', radius=0, shape=MSO_SHAPE.RECTANGLE)
for i,(yy,cc) in enumerate([(2.3,'sage7'),(2.53,'sage6'),(2.76,'echon2'),(2.99,'echon'),(3.22,'gold')]):
    box(s, 0.7, yy, 2.0, 0.16, fill=cc, line=None, radius=0.3)
text(s, "Win the carpenter,\narm the counter,\nbuild the brand.", 0.7, 3.7, 9, 2.0, size=40, color='white', font=TITLE, bold=True, line_spacing=1.06, track=-30)
text(s, "A sales-facilitation engine, reported monthly — every rupee tied to an outcome.", 0.72, 5.95, 10, 0.4, size=13.5, color='echonl', font=BODY)
text(s, "business@sagemedia.in", 0.72, 6.45, 6, 0.4, size=14, color='white', font=BODY, bold=True)
text(s, "Sage Media · Jaipur · Mumbai · Dubai", 0.72, 6.84, 7, 0.4, size=11, color='cream3', font=BODY)
text(s, "Echon", 11.0, 0.7, 1.8, 0.6, size=24, color='echonl', font=TITLE, bold=True, align=PP_ALIGN.RIGHT)
text(s, "Annual Marketing Plan\nBudget & Return · Year 1 · India", 8.5, 1.35, 4.3, 0.8, size=11, color='cream3', font=BODY, align=PP_ALIGN.RIGHT, line_spacing=1.15)

out='/home/user/Sage-portfolio/Echon-Annual-Marketing-Plan-SageMedia.pptx'
prs.save(out)
print("saved", out)
