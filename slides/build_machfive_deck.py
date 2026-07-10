#!/usr/bin/env python3
"""Mach Five — Go-to-Market & Growth Plan (Sage Media).
Elevated editable pitch deck on our format + Glacial Indifference/Inter.
Functional health drink (creatine, vitamin C — energy, strength, recovery, brain)."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

C = {
    'paper':'F3F1E8','paper2':'E9E5D7','paper3':'DED9C7',
    'ink':'0F1216','ink2':'181C24','ink3':'262B35',
    'volt':'5FA80E','voltL':'A6D93B','voltD':'2E5406',   # energy / strength / recovery
    'focus':'2B5BFF','focusL':'6E8CFF',                    # brain / focus
    'citrus':'F2A100',                                     # vitamin C pop
    'mute':'6B7168','muteL':'A9AEA0','line':'D5D0C1','white':'FFFFFF',
}
def rgb(h): return RGBColor.from_string(C.get(h,h))
TITLE='Glacial Indifference'; BODY='Inter'
prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
BLANK=prs.slide_layouts[6]; SW,SH=prs.slide_width,prs.slide_height

def slide(bg='paper'):
    s=prs.slides.add_slide(BLANK)
    r=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,0,0,SW,SH)
    r.fill.solid(); r.fill.fore_color.rgb=rgb(bg); r.line.fill.background(); r.shadow.inherit=False
    return s
def _font(run,size,color,font,bold,italic=False,spc=None):
    run.font.size=Pt(size); run.font.bold=bold; run.font.italic=italic
    run.font.name=font; run.font.color.rgb=rgb(color)
    rPr=run._r.get_or_add_rPr()
    if spc is not None: rPr.set('spc',str(int(spc)))
    for tag in ('a:latin','a:cs'):
        e=rPr.find(qn(tag))
        if e is None: e=rPr.makeelement(qn(tag),{}); rPr.append(e)
        e.set('typeface',font)
def text(s,txt,l,t,w,h,size=14,color='ink',font=BODY,bold=False,align=PP_ALIGN.LEFT,
         anchor=MSO_ANCHOR.TOP,italic=False,ls=1.05,track=None,sp_after=2):
    tb=s.shapes.add_textbox(Inches(l),Inches(t),Inches(w),Inches(h))
    tf=tb.text_frame; tf.word_wrap=True; tf.vertical_anchor=anchor
    tf.margin_left=0;tf.margin_right=0;tf.margin_top=0;tf.margin_bottom=0
    for i,ln in enumerate(txt.split('\n')):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.alignment=align; p.space_after=Pt(sp_after); p.space_before=Pt(0)
        try: p.line_spacing=ls
        except Exception: pass
        r=p.add_run(); r.text=ln; _font(r,size,color,font,bold,italic,spc=track)
    return tb
def box(s,l,t,w,h,fill='white',line=None,line_w=1.0,radius=0.08,shape=MSO_SHAPE.ROUNDED_RECTANGLE,shadow=False):
    sp=s.shapes.add_shape(shape,Inches(l),Inches(t),Inches(w),Inches(h))
    if fill is None: sp.fill.background()
    else: sp.fill.solid(); sp.fill.fore_color.rgb=rgb(fill)
    if line is None: sp.line.fill.background()
    else: sp.line.color.rgb=rgb(line); sp.line.width=Pt(line_w)
    sp.shadow.inherit=False
    if shape==MSO_SHAPE.ROUNDED_RECTANGLE:
        try: sp.adjustments[0]=radius
        except Exception: pass
    if shadow:
        el=sp._element.spPr; ef=el.makeelement(qn('a:effectLst'),{})
        sh=el.makeelement(qn('a:outerShdw'),{'blurRad':'150000','dist':'50000','dir':'5400000','rotWithShape':'0'})
        clr=el.makeelement(qn('a:srgbClr'),{'val':'0F1216'}); al=el.makeelement(qn('a:alpha'),{'val':'18000'})
        clr.append(al); sh.append(clr); ef.append(sh); el.append(ef)
    return sp
def grad(sp,c1,c2,angle=55):
    spPr=sp._element.spPr
    for tag in ('a:solidFill','a:noFill','a:gradFill','a:blipFill','a:pattFill'):
        e=spPr.find(qn(tag));  spPr.remove(e) if e is not None else None
    g=spPr.makeelement(qn('a:gradFill'),{}); lst=spPr.makeelement(qn('a:gsLst'),{})
    for pos,c in [(0,c1),(100000,c2)]:
        gs=spPr.makeelement(qn('a:gs'),{'pos':str(pos)}); clr=spPr.makeelement(qn('a:srgbClr'),{'val':C.get(c,c)})
        gs.append(clr); lst.append(gs)
    g.append(lst); g.append(spPr.makeelement(qn('a:lin'),{'ang':str(int(angle*60000)),'scaled':'1'}))
    ln=spPr.find(qn('a:ln'))
    spPr.insert(list(spPr).index(ln),g) if ln is not None else spPr.append(g)
    return sp
def shape_text(sp,txt,size=14,color='white',font=BODY,bold=False,align=PP_ALIGN.LEFT,anchor=MSO_ANCHOR.MIDDLE,track=None):
    tf=sp.text_frame; tf.word_wrap=True; tf.vertical_anchor=anchor
    tf.margin_left=Inches(0.12);tf.margin_right=Inches(0.12);tf.margin_top=Inches(0.03);tf.margin_bottom=Inches(0.03)
    for i,ln in enumerate(txt.split('\n')):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph(); p.alignment=align
        r=p.add_run(); r.text=ln; _font(r,size,color,font,bold,spc=track)
    return sp
def imgblock(s,l,t,w,h,c1='ink2',c2='volt',label='IMAGERY',angle=55,radius=0.05):
    b=box(s,l,t,w,h,fill='ink2',radius=radius); grad(b,c1,c2,angle)
    cx=l+w/2; cy=t+h/2
    o=box(s,cx-0.33,cy-0.33,0.66,0.66,fill=None,line='white',line_w=1.6,radius=0.5,shape=MSO_SHAPE.OVAL); o.fill.background()
    box(s,cx-0.08,cy-0.08,0.16,0.16,fill='white',radius=0.5,shape=MSO_SHAPE.OVAL)
    chip=box(s,l+0.22,t+h-0.56,0.09+0.075*len(label),0.34,fill='ink',radius=0.5); chip.fill.fore_color.rgb=rgb('0F1216')
    shape_text(chip,label,size=8,color='white',bold=True,align=PP_ALIGN.CENTER,track=40)
    return b
def eyebrow(s,txt,l=0.85,t=0.7,dark=False,color=None):
    return text(s,txt.upper(),l,t,9,0.35,size=10.5,color=(color or ('voltL' if dark else 'volt')),font=BODY,bold=True,track=200)
def footer(s,idx,total=20,dark=False):
    col='muteL' if dark else 'mute'
    box(s,0.85,7.03,11.63,0.014,fill=('ink3' if dark else 'line'),radius=0,shape=MSO_SHAPE.RECTANGLE)
    text(s,"SAGE MEDIA  ·  MACH FIVE",0.85,7.1,6,0.3,size=8.5,color=col,font=BODY,bold=True,track=100)
    text(s,f"{idx:02d} / {total:02d}",11.0,7.1,1.48,0.3,size=8.5,color=col,font=BODY,align=PP_ALIGN.RIGHT,track=80)
def sectnum(s,n,color='ink3'):
    text(s,n,8.0,-0.7,6,6,size=300,color=color,font=TITLE,bold=True,align=PP_ALIGN.LEFT,track=-60)
def divider(s,num,label,title,sub,idx):
    sectnum(s,num,'ink3'); box(s,0.85,0,0.22,SH,fill='volt',radius=0,shape=MSO_SHAPE.RECTANGLE)
    text(s,label,1.35,2.35,7,0.4,size=11,color='voltL',font=BODY,bold=True,track=220)
    text(s,title,1.3,2.82,8.2,2.2,size=58,color='white',font=TITLE,bold=True,track=-40,ls=0.98)
    text(s,sub,1.35,5.25,7.6,0.9,size=14.5,color='muteL',font=BODY,ls=1.3)
    footer(s,idx,dark=True)
def header(s,eyb,title,sub=None,tsize=30):
    box(s,0,0,SW,Inches(0.12),fill='volt',radius=0,shape=MSO_SHAPE.RECTANGLE)
    eyebrow(s,eyb)
    text(s,title,0.85,1.12,11.8,1.1,size=tsize,color='ink',font=TITLE,bold=True,track=-30,ls=1.0)
    if sub: text(s,sub,0.85,1.12+(0.62 if '\n' not in title else 1.15),11.4,0.7,size=12.5,color='mute',font=BODY,ls=1.2)

# horizontal bars inside a region
def hbars(s,rows,x,y,label_w,full_w,bh=0.34,gap=0.22,maxv=None,valsuffix=''):
    maxv=maxv or max(r[1] for r in rows)
    yy=y
    for nm,v,col in rows:
        text(s,nm,x,yy-0.02,label_w-0.15,0.35,size=10.5,color='ink',font=BODY,bold=True,align=PP_ALIGN.RIGHT,anchor=MSO_ANCHOR.MIDDLE)
        box(s,x+label_w,yy,full_w,bh,fill='paper3',radius=0.4)
        box(s,x+label_w,yy,max(full_w*v/maxv,0.12),bh,fill=col,radius=0.4)
        text(s,f"{v}{valsuffix}",x+label_w+full_w*v/maxv+0.1,yy-0.02,1.2,0.35,size=10.5,color='ink',font=TITLE,bold=True,anchor=MSO_ANCHOR.MIDDLE)
        yy+=bh+gap
    return yy

# vertical columns
def columns(s,rows,x,base_y,col_w,gap,maxh,maxv=None,colr='volt',lblcolor='mute',valcolor='ink'):
    maxv=maxv or max(r[1] for r in rows)
    xx=x
    for nm,v in rows:
        hgt=maxh*v/maxv
        box(s,xx,base_y-hgt,col_w,hgt,fill=colr,radius=0.14)
        text(s,str(v),xx-0.2,base_y-hgt-0.3,col_w+0.4,0.28,size=9.5,color=valcolor,font=TITLE,bold=True,align=PP_ALIGN.CENTER)
        text(s,nm,xx-0.2,base_y+0.06,col_w+0.4,0.3,size=9,color=lblcolor,font=BODY,bold=True,align=PP_ALIGN.CENTER)
        xx+=col_w+gap
    return xx

# ============================================================= 1 COVER
s=slide('ink')
box(s,0,0,0.22,SH,fill='volt',radius=0,shape=MSO_SHAPE.RECTANGLE)
# right imagery
imgblock(s,8.55,0,4.78,4.35,'ink2','volt','PRODUCT HERO',angle=60,radius=0.0)
imgblock(s,8.55,4.35,4.78,3.15,'focus','ink2','STRENGTH · FOCUS',angle=35,radius=0.0)
text(s,"GO-TO-MARKET & GROWTH PLAN · 2026",0.95,0.9,7,0.4,size=11,color='voltL',font=BODY,bold=True,track=200)
text(s,"Mach Five",0.9,1.75,8,1.4,size=74,color='white',font=TITLE,bold=True,track=-55)
box(s,0.98,3.25,0.9,0.09,fill='volt',radius=0,shape=MSO_SHAPE.RECTANGLE)
text(s,"Creatine · Vitamin C · functional energy —\nfor strength, recovery and a sharper mind.",
     0.95,3.55,6.9,1.0,size=15.5,color='muteL',font=BODY,ls=1.3)
text(s,"Prepared for  ·  Aryaman Singh Shaktawat",0.95,5.75,6.5,0.35,size=11,color='mute',font=BODY,bold=True,track=30)
text(s,"Prepared by Sage Media",0.95,6.55,6,0.35,size=12.5,color='white',font=TITLE,bold=True)
text(s,"Jaipur · Mumbai · Dubai  ·  business@sagemedia.in",0.95,6.95,7.5,0.3,size=9.5,color='mute',font=BODY,track=20)

# ============================================================= 2 EXEC SUMMARY
s=slide('paper')
box(s,0,0,SW,Inches(0.12),fill='volt',radius=0,shape=MSO_SHAPE.RECTANGLE)
eyebrow(s,"Executive Summary")
text(s,"A launch built on four compounding tailwinds.",0.85,1.1,11.6,0.9,size=29,color='ink',font=TITLE,bold=True,track=-30)
text(s,"Mach Five enters at the intersection of India's supplement boom and the global creatine + brain-health surge. This is a D2C-first plan: build the brand, prove the funnel, then reinvest behind what sells.",
     0.85,1.85,11.2,0.8,size=13,color='mute',font=BODY,ls=1.3)
rows=[("Recommended entry","D2C-first (own store) → marketplaces → quick-commerce → offline, sequenced to sales."),
      ("Marketing investment","₹3–11 L / month (ramped) + one-time brand & store build; scaled behind proven ROAS."),
      ("Turnaround","First traction by Month 3 · contribution breakeven ~Month 6 · profitable scale by Month 9–12."),
      ("Year-1 ambition","Build to a ₹40 L+ monthly run-rate; ~₹4–5 Cr Year-1 revenue enabled (illustrative)."),
      ("Reporting","Monthly — against CAC, ROAS, repeat rate and run-rate, not vanity metrics.")]
yy=3.0
for k,v in rows:
    box(s,0.85,yy+0.5,8.0,0.012,fill='line',radius=0,shape=MSO_SHAPE.RECTANGLE)
    text(s,k.upper(),0.85,yy,2.6,0.4,size=10,color='volt',font=BODY,bold=True,track=60,anchor=MSO_ANCHOR.MIDDLE)
    text(s,v,3.55,yy,5.3,0.5,size=11.5,color='ink',font=BODY,ls=1.1,anchor=MSO_ANCHOR.MIDDLE)
    yy+=0.68
# highlight card
hc=box(s,9.15,2.9,3.35,3.6,fill='ink',radius=0.06,shadow=True)
text(s,"THE CATEGORY, COMPOUNDING",9.45,3.2,2.8,0.35,size=9.5,color='voltL',font=BODY,bold=True,track=120)
text(s,"25%",9.45,3.6,2.8,1.0,size=58,color='white',font=TITLE,bold=True,track=-30)
text(s,"global creatine-supplement CAGR — the fastest-growing corner of a booming market.",
     9.45,4.75,2.75,1.0,size=12,color='muteL',font=BODY,ls=1.25)
box(s,9.45,5.95,2.75,0.44,fill='ink3',radius=0.4)
box(s,9.45,5.95,2.75*0.72,0.44,fill='volt',radius=0.4)
text(s,"India, still early",9.55,5.98,2.5,0.38,size=9.5,color='white',font=BODY,bold=True,anchor=MSO_ANCHOR.MIDDLE)
footer(s,2)

# ============================================================= 3 DIVIDER — MARKET
divider(s:=slide('ink'),"01","THE MARKET","The\nOpportunity","Market size, growth and the white space Mach Five can own.",3)

# ============================================================= 4 MARKET TAILWINDS
s=slide('paper')
header(s,"The Market · Tailwinds","Mach Five rides the fastest-growing categories in nutrition.",
       "5-year CAGR by category — the brand sits where growth is compounding hardest.")
box(s,0.85,2.6,7.5,3.95,fill='white',line='line',line_w=1,radius=0.05,shadow=True)
text(s,"Category growth · CAGR to 2030",1.2,2.85,6,0.4,size=12.5,color='ink',font=TITLE,bold=True)
bars=[("Creatine (global)",25.0,'volt'),("Nootropics (global)",14.7,'focus'),
      ("Brain-health supp.",13.7,'focus'),("India supplements",13.1,'volt'),("India functional drinks",11.7,'volt')]
hbars(s,bars,1.2,3.45,2.35,3.4,bh=0.34,gap=0.24,maxv=25.0,valsuffix='%')
# right facts
facts=[("$5.2B","India dietary-supplements market (2024) — growing 13% a year.",'volt'),
       ("$3.8B → $7.4B","India functional-drinks market, 2024 → 2030.",'focus'),
       ("$1.4B → $4.2B","Global creatine supplements, 2025 → 2030.",'citrus')]
fy=2.6
for v,l,col in facts:
    box(s,8.65,fy,3.85,1.2,fill='paper2',line='line',line_w=1,radius=0.09)
    box(s,8.65,fy,0.1,1.2,fill=col,radius=0,shape=MSO_SHAPE.RECTANGLE)
    text(s,v,8.95,fy+0.14,3.5,0.5,size=23,color='ink',font=TITLE,bold=True,track=-20)
    text(s,l,8.95,fy+0.66,3.4,0.5,size=10,color='mute',font=BODY,ls=1.15)
    fy+=1.33
footer(s,4)

# ============================================================= 5 MARKET FORECAST
s=slide('paper2')
header(s,"The Market · Forecast","India's supplement market roughly doubles by 2030.",
       "India dietary-supplements market, USD billion — a rising tide Mach Five launches into.")
box(s,0.85,2.55,11.63,4.0,fill='paper',line='line',line_w=1,radius=0.05,shadow=True)
yrs=[("2024",5.2),("2025",5.9),("2026",6.6),("2027",7.5),("2028",8.5),("2029",9.5),("2030",10.8)]
columns(s,yrs,1.45,5.85,0.82,0.28,3.0,maxv=10.8,colr='volt')
text(s,"USD billion",1.1,3.0,2,0.3,size=10,color='mute',font=BODY,bold=True,track=40)
# CAGR callout (right gutter, no overlap)
cc=box(s,9.55,3.35,2.7,1.6,fill='ink',radius=0.08,shadow=True)
text(s,"13.1%",9.8,3.6,2.3,0.8,size=42,color='voltL',font=TITLE,bold=True,track=-25)
text(s,"CAGR · 2025–2030",9.8,4.5,2.3,0.35,size=10.5,color='muteL',font=BODY,bold=True,track=40)
text(s,"Sources: Grand View Research; MarketsandData; Market.us. Figures illustrative and directional.",
     0.85,6.7,11,0.3,size=8.5,color='mute',font=BODY,italic=True)
footer(s,5)

# ============================================================= 6 TAM SAM SOM
s=slide('paper')
header(s,"The Market · Sizing","Where Mach Five plays — and what it can win.")
# concentric ovals
cx,cy=3.9,4.4
for (w,h,col,lbl) in [(5.6,4.2,'paper3','TAM'),(3.9,2.95,'voltL','SAM'),(2.2,1.65,'volt','SOM')]:
    o=box(s,cx-w/2,cy-h/2,w,h,fill=col,line=None,radius=0.5,shape=MSO_SHAPE.OVAL)
text(s,"TAM",cx-0.6,cy-1.95,1.2,0.3,size=11,color='mute',font=BODY,bold=True,align=PP_ALIGN.CENTER,track=80)
text(s,"SOM",cx-0.6,cy-0.15,1.2,0.3,size=11,color='white',font=BODY,bold=True,align=PP_ALIGN.CENTER,track=60)
# legend right
items=[("TAM · ₹75,000 Cr","India supplements + functional beverages — the total category.",'paper3'),
       ("SAM · ₹5,000 Cr","Premium, online-first performance & cognition nutrition (est.).",'voltL'),
       ("SOM · ₹15–20 Cr","Mach Five's realistic 3-year D2C capture (est.).",'volt')]
iy=2.7
for v,l,col in items:
    box(s,7.5,iy,0.34,0.34,fill=col,line='line',line_w=1,radius=0.18)
    text(s,v,7.98,iy-0.05,4.6,0.4,size=15,color='ink',font=TITLE,bold=True)
    text(s,l,7.98,iy+0.42,4.5,0.6,size=10.5,color='mute',font=BODY,ls=1.2)
    iy+=1.25
text(s,"The wedge: win a focused, winnable slice online first — then expand the addressable market channel by channel.",
     7.5,6.35,4.9,0.6,size=10.5,color='volt',font=BODY,bold=True,ls=1.2)
footer(s,6)

# ============================================================= 7 DIVIDER — BRAND
divider(s:=slide('ink'),"02","BRAND & POSITIONING","The\nBrand","How Mach Five shows up, who it's for, and what it looks like.",7)

# ============================================================= 8 POSITIONING
s=slide('paper')
header(s,"Brand · Positioning","One drink, four jobs.",
       "Mach Five is the functional daily fuel for people who train hard and think harder.")
pill=[("Energy",'volt',"Clean, sustained energy — no crash."),
      ("Strength",'volt',"Creatine-led performance & power output."),
      ("Recovery",'citrus',"Vitamin C + actives for faster bounce-back."),
      ("Brain",'focus',"Focus and cognition for work and study.")]
cw=2.86; gap=0.22; x0=0.85; y0=2.75; ch=3.4
for i,(h,col,b) in enumerate(pill):
    x=x0+i*(cw+gap)
    box(s,x,y0,cw,ch,fill='white',line='line',line_w=1,radius=0.06,shadow=True)
    box(s,x,y0,cw,0.1,fill=col,radius=0,shape=MSO_SHAPE.RECTANGLE)
    text(s,f"0{i+1}",x+0.28,y0+0.35,cw-0.5,0.7,size=34,color='paper3',font=TITLE,bold=True)
    text(s,h,x+0.28,y0+1.35,cw-0.5,0.5,size=20,color='ink',font=TITLE,bold=True)
    text(s,b,x+0.28,y0+1.95,cw-0.56,1.0,size=11.5,color='mute',font=BODY,ls=1.25)
text(s,"Positioning line:  “Fuel for the fast — energy, strength, recovery and focus in one clean daily hit.”",
     0.85,6.45,11.6,0.4,size=12,color='volt',font=BODY,bold=True)
footer(s,8)

# ============================================================= 9 COMPETITIVE MAP
s=slide('paper2')
header(s,"Brand · Landscape","A clear gap: premium, brain-and-body, made for India.")
# quadrant
qx,qy,qs=1.6,2.7,3.7
box(s,qx,qy,7.2,qs,fill='paper',line='line',line_w=1,radius=0.04,shadow=True)
box(s,qx,qy+qs/2,7.2,0.012,fill='line',radius=0,shape=MSO_SHAPE.RECTANGLE)
box(s,qx+7.2/2,qy,0.012,qs,fill='line',radius=0,shape=MSO_SHAPE.RECTANGLE)
text(s,"PREMIUM",qx+7.2/2-1,qy-0.32,2,0.3,size=9,color='mute',font=BODY,bold=True,align=PP_ALIGN.CENTER,track=80)
text(s,"MASS",qx+7.2/2-1,qy+qs+0.06,2,0.3,size=9,color='mute',font=BODY,bold=True,align=PP_ALIGN.CENTER,track=80)
text(s,"BODY ONLY",qx-1.15,qy+qs/2-0.15,1.4,0.3,size=9,color='mute',font=BODY,bold=True,align=PP_ALIGN.CENTER,track=40)
text(s,"BODY + BRAIN",qx+7.2-0.35,qy+qs/2-0.15,1.7,0.3,size=9,color='mute',font=BODY,bold=True,align=PP_ALIGN.LEFT,track=40)
dots=[("Legacy protein",qx+1.3,qy+2.7,'muteL'),("Energy drinks",qx+1.9,qy+3.0,'muteL'),
      ("Global creatine",qx+2.4,qy+1.1,'muteL'),("Nootropic startups",qx+5.4,qy+1.5,'muteL')]
for nm,dx,dy,col in dots:
    box(s,dx-0.09,dy-0.09,0.18,0.18,fill=col,radius=0.5,shape=MSO_SHAPE.OVAL)
    text(s,nm,dx+0.16,dy-0.14,2.1,0.3,size=9,color='mute',font=BODY)
# Mach Five sweet spot (top-right)
mx,my=qx+5.6,qy+0.95
box(s,mx-0.16,my-0.16,0.32,0.32,fill='volt',radius=0.5,shape=MSO_SHAPE.OVAL)
text(s,"MACH FIVE",mx+0.24,my-0.2,2.2,0.4,size=11,color='volt',font=TITLE,bold=True)
# right note
box(s,9.15,2.7,3.35,3.7,fill='ink',radius=0.06,shadow=True)
text(s,"THE WHITE SPACE",9.45,2.95,2.8,0.35,size=9.5,color='voltL',font=BODY,bold=True,track=120)
for i,t in enumerate(["Premium, not mass.","Body + brain in one.","Made for Indian routines & taste.","D2C storytelling, not shelf clutter."]):
    text(s,"—  "+t,9.45,3.5+i*0.62,2.9,0.5,size=12.5,color='white',font=BODY,ls=1.15)
footer(s,9)

# ============================================================= 10 PERSONAS
s=slide('paper')
header(s,"Brand · Audience","Three people we build for.")
ps=[("The Lifter","22–32 · gym 4–6×/wk","Wants clean creatine & recovery that actually works. Reads labels, follows fitness creators.",'volt'),
    ("The Grinder","24–35 · founder / knowledge worker","Needs energy + focus without the crash. Buys into performance & brain health.",'focus'),
    ("The Optimiser","20–30 · student / early-career","Health-curious, values-driven, discovers brands on Reels; wants aspirational + affordable.",'citrus')]
cw=3.86; gap=0.22; x0=0.85; y0=2.6; ch=3.9
for i,(h,meta,b,col) in enumerate(ps):
    x=x0+i*(cw+gap)
    box(s,x,y0,cw,ch,fill='white',line='line',line_w=1,radius=0.06,shadow=True)
    imgblock(s,x+0.28,y0+0.3,cw-0.56,1.5,'ink2',col,'PERSONA',angle=45,radius=0.08)
    text(s,h,x+0.3,y0+1.95,cw-0.6,0.4,size=18,color='ink',font=TITLE,bold=True)
    text(s,meta.upper(),x+0.3,y0+2.42,cw-0.6,0.3,size=9,color=col,font=BODY,bold=True,track=40)
    text(s,b,x+0.3,y0+2.78,cw-0.6,1.0,size=11,color='mute',font=BODY,ls=1.22)
footer(s,10)

# ============================================================= 11 MOOD BOARD
s=slide('ink')
box(s,0,0,SW,Inches(0.12),fill='volt',radius=0,shape=MSO_SHAPE.RECTANGLE)
text(s,"BRAND · VISUAL REFERENCE",0.85,0.55,9,0.35,size=10.5,color='voltL',font=BODY,bold=True,track=200)
text(s,"Mood board & art direction.",0.85,0.95,11,0.8,size=28,color='white',font=TITLE,bold=True,track=-30)
# imagery tiles
tiles=[("PRODUCT HERO",'ink2','volt'),("STRENGTH / GYM",'ink2','focus'),("CITRUS / VIT-C",'ink2','citrus'),
       ("BRAIN / FOCUS",'focus','ink'),("RECOVERY / CALM",'ink2','volt'),("LIFESTYLE",'ink2','citrus')]
tx=0.85; tw=2.62; th=1.75; gapx=0.16
for i,(lbl,a,b) in enumerate(tiles):
    col=i%4; row=i//4
    if i<4: x=0.85+i*(tw+gapx); y=1.95
    else: x=0.85+(i-4)*(tw+gapx); y=3.86
    imgblock(s,x,y,tw,th,a,b,lbl,angle=45+i*10,radius=0.05)
# palette + type panel (right of second row)
px=6.2; py=3.86
box(s,px,py,6.28,1.75,fill='ink2',radius=0.06)
text(s,"PALETTE",px+0.3,py+0.22,3,0.3,size=9,color='muteL',font=BODY,bold=True,track=120)
for i,(c,nm) in enumerate([('volt','Volt'),('focus','Focus'),('citrus','Vit-C'),('ink','Ink'),('paper','Paper')]):
    sw=box(s,px+0.3+i*0.66,py+0.55,0.52,0.52,fill=c,line='ink3',line_w=1,radius=0.16)
    text(s,nm,px+0.28+i*0.66,py+1.12,0.62,0.25,size=7.5,color='muteL',font=BODY,align=PP_ALIGN.CENTER)
text(s,"TYPE  ·  Glacial Indifference / Inter",px+3.9,py+0.22,3,0.3,size=9,color='muteL',font=BODY,bold=True,track=60)
text(s,"Aa",px+3.9,py+0.5,2,0.7,size=40,color='white',font=TITLE,bold=True)
text(s,"bold display · clean body",px+5.0,py+0.72,1.3,0.6,size=8.5,color='muteL',font=BODY,ls=1.1)
text(s,"Direction: matte-black + volt-green, high-contrast studio product shots, real Indian athletes, science-forward brain motifs, citrus macro pops. Full AI-image prompts in the appendix.",
     0.85,5.85,11.6,0.7,size=10.5,color='muteL',font=BODY,ls=1.3)
footer(s,11,dark=True)

# ============================================================= 12 DIVIDER — GTM
divider(s:=slide('ink'),"03","GO-TO-MARKET","Market\nEntry","The wedge, the channels and the 12-month roadmap.",12)

# ============================================================= 13 ENTRY STRATEGY
s=slide('paper')
header(s,"Go-to-Market · The Wedge","Enter narrow, expand by channel.",
       "Win D2C first with a hero SKU, then unlock each channel as sales prove out.")
steps=[("01","D2C store","Own the story & margins","Shopify hero-SKU launch, content + creators, performance ads.",'volt'),
       ("02","Marketplaces","Capture demand","Amazon & Flipkart once reviews & ROAS are proven.",'focus'),
       ("03","Quick-commerce","Win impulse & repeat","Blinkit / Zepto / Instamart in metros at scale.",'citrus'),
       ("04","Offline & gyms","Trust & trial","Gym tie-ups and select retail as the brand compounds.",'volt')]
cw=2.86; gap=0.22; x0=0.85; y0=2.75; ch=3.5
for i,(n,h,tag,b,col) in enumerate(steps):
    x=x0+i*(cw+gap)
    box(s,x,y0,cw,ch,fill='white',line='line',line_w=1,radius=0.06,shadow=True)
    text(s,n,x+0.28,y0+0.3,cw-0.5,0.7,size=30,color=col,font=TITLE,bold=True)
    text(s,h,x+0.28,y0+1.15,cw-0.5,0.45,size=17,color='ink',font=TITLE,bold=True)
    text(s,tag.upper(),x+0.28,y0+1.62,cw-0.5,0.3,size=8.5,color=col,font=BODY,bold=True,track=50)
    text(s,b,x+0.28,y0+2.0,cw-0.56,1.1,size=10.5,color='mute',font=BODY,ls=1.22)
    if i<3: text(s,"→",x+cw-0.02,y0+1.25,0.4,0.4,size=18,color='paper3',font=BODY,bold=True)
footer(s,13)

# ============================================================= 14 ROADMAP
s=slide('paper2')
header(s,"Go-to-Market · Roadmap","A phased 12-month launch — with a clear turnaround point.")
ph=[("Q1 · M1–3","Build & Validate","Brand, Shopify store, hero SKU launch, seed content, micro-influencers, creative testing.","→ First ₹5 L MRR",'volt'),
    ("Q2 · M4–6","Traction & Breakeven","Scale winning ads, Amazon launch, subscriptions, UGC engine. Contribution breakeven ~M6.","→ ₹15 L MRR",'focus'),
    ("Q3 · M7–9","Scale","Quick-commerce, SKU 2, brand ambassador, PR push, retention flows.","→ ₹30 L MRR",'citrus'),
    ("Q4 · M10–12","Expand","Offline & gym pilots, SKU 3, category push, always-on performance.","→ ₹40 L+ MRR",'volt')]
cw=2.86; gap=0.22; x0=0.85; y0=2.5; ch=4.0
for i,(q,h,b,mile,col) in enumerate(ph):
    x=x0+i*(cw+gap); dark=(i==1)
    box(s,x,y0,cw,ch,fill=('ink' if dark else 'paper'),line=(None if dark else 'line'),line_w=1,radius=0.05,shadow=True)
    text(s,q.upper(),x+0.3,y0+0.35,cw-0.6,0.3,size=10,color=('voltL' if dark else col),font=BODY,bold=True,track=80)
    text(s,h,x+0.3,y0+0.72,cw-0.6,0.5,size=18,color=('white' if dark else 'ink'),font=TITLE,bold=True)
    text(s,b,x+0.3,y0+1.4,cw-0.6,1.7,size=10.5,color=('muteL' if dark else 'mute'),font=BODY,ls=1.25)
    box(s,x+0.3,y0+ch-0.85,cw-0.6,0.55,fill=(col if not dark else 'ink3'),radius=0.3)
    shape_text(box(s,x+0.3,y0+ch-0.85,cw-0.6,0.55,fill=col,radius=0.3),mile,size=11,color=('ink' if col in('citrus','voltL','volt') else 'white'),bold=True,align=PP_ALIGN.CENTER)
footer(s,14)

# ============================================================= 15 DIVIDER — SPEND & RETURN
divider(s:=slide('ink'),"04","BUDGET & RETURN","The\nNumbers","Where the money goes, the milestone ladder, and the payback.",15)

# ============================================================= 16 BUDGET ALLOCATION
s=slide('paper')
header(s,"Budget · Monthly Spend","Where the marketing budget goes at scale.",
       "Illustrative monthly marketing mix once the engine is proven (~₹11 L). Media-led, content-fuelled.")
# stacked 100% bar
segs=[("Performance media",45,'volt'),("Content & production",18,'focus'),("Influencer & UGC",14,'citrus'),
      ("Founder branding",10,'ink2'),("PR & earned",5,'muteL'),("Retainer & tools",8,'paper3')]
bx=0.85; bw=11.63; by=2.75; bh=0.7; acc=0
for nm,pct,col in segs:
    w=bw*pct/100
    seg=box(s,bx+acc,by,w,bh,fill=col,line='paper',line_w=1.5,radius=0,shape=MSO_SHAPE.RECTANGLE)
    shape_text(seg,f"{pct}%",size=12,color=('ink' if col in('citrus','paper3','muteL') else 'white'),bold=True,align=PP_ALIGN.CENTER)
    acc+=w
lx=0.85; ly=3.7
for nm,pct,col in segs:
    box(s,lx,ly+0.04,0.2,0.2,fill=col,line='line',line_w=0.75,radius=0.2)
    text(s,nm,lx+0.28,ly,2.2,0.4,size=9.5,color='ink',font=BODY,ls=1.0)
    lx+=0.28+0.075*len(nm)+0.42
# one-time + phasing note cards
cards=[("₹50 K","One-time · Branding & identity"),("₹60 K","One-time · Shopify store build"),
       ("₹3 → 11 L","Monthly marketing, ramped by proof"),("45%","Goes to performance media at scale")]
cy2=4.7; cwid=2.86; g2=0.22
for i,(v,l) in enumerate(cards):
    x=0.85+i*(cwid+g2)
    box(s,x,cy2,cwid,1.35,fill='paper2',line='line',line_w=1,radius=0.08)
    text(s,v,x+0.28,cy2+0.2,cwid-0.5,0.6,size=25,color='ink',font=TITLE,bold=True,track=-20)
    text(s,l,x+0.28,cy2+0.78,cwid-0.5,0.5,size=9.5,color='mute',font=BODY,ls=1.15)
text(s,"Excludes Shopify/app subscriptions, shoot actuals (models, travel, food, studio) and GST — billed per the deliverables proposal.",
     0.85,6.35,11.6,0.3,size=8.5,color='mute',font=BODY,italic=True)
footer(s,16)

# ============================================================= 17 MILESTONE LADDER
s=slide('paper2')
header(s,"Growth · Milestone Ladder","What we do after each sales milestone.",
       "A reinvestment ladder — every rupee of new revenue unlocks the next channel and lever.")
rungs=[("LAUNCH → ₹5 L MRR","Prove the funnel","Hero SKU · D2C store · micro-influencers · creative testing · seed reviews.",'paper','volt'),
       ("₹5 L → ₹15 L MRR","Scale what works","Scale winning ads · Amazon + Flipkart · subscriptions · UGC engine.",'paper','focus'),
       ("₹15 L → ₹40 L MRR","Add channels","Quick-commerce (Blinkit/Zepto) · SKU 2–3 · brand ambassador · PR.",'ink','citrus'),
       ("₹40 L+ MRR","Go omnichannel","Offline & gym retail · category expansion · always-on performance.",'ink','volt')]
n=len(rungs); baseY=6.4; x0=0.95; step_w=2.75; gap=0.28; hbase=1.4
for i,(head,sub,b,fill,col) in enumerate(rungs):
    x=x0+i*(step_w+gap); h=hbase+i*0.95; y=baseY-h
    dark=(fill=='ink')
    box(s,x,y,step_w,h,fill=fill,line=(None if dark else 'line'),line_w=1,radius=0.05,shadow=True)
    box(s,x,y,step_w,0.09,fill=col,radius=0,shape=MSO_SHAPE.RECTANGLE)
    text(s,head,x+0.24,y+0.24,step_w-0.45,0.5,size=12.5,color=(col if not dark else 'voltL'),font=TITLE,bold=True)
    text(s,sub.upper(),x+0.24,y+0.74,step_w-0.45,0.3,size=8.5,color=('muteL' if dark else 'mute'),font=BODY,bold=True,track=50)
    text(s,b,x+0.24,y+1.08,step_w-0.45,h-1.2,size=9.5,color=('muteL' if dark else 'mute'),font=BODY,ls=1.22)
footer(s,17)

# ============================================================= 18 RETURNS / TURNAROUND
s=slide('paper')
header(s,"Return · Turnaround","Revenue ramps; contribution turns positive by Month 6.",
       "Illustrative monthly revenue (₹ Lakh). Every input is tuned to Mach Five's real numbers at kickoff.")
box(s,0.85,2.55,8.0,4.0,fill='white',line='line',line_w=1,radius=0.05,shadow=True)
mrev=[("M1",3),("M2",6),("M3",10),("M4",15),("M5",22),("M6",30),("M7",38),("M8",47),("M9",56),("M10",66),("M11",78),("M12",90)]
# columns
bx0=1.35; cwid=0.5; g=0.12; maxh=2.7; base=5.9; maxv=90
xx=bx0
for i,(nm,v) in enumerate(mrev):
    hgt=maxh*v/maxv; col='volt' if i>=5 else 'paper3'
    box(s,xx,base-hgt,cwid,hgt,fill=col,radius=0.12)
    text(s,nm,xx-0.1,base+0.06,cwid+0.2,0.25,size=7.5,color='mute',font=BODY,align=PP_ALIGN.CENTER)
    xx+=cwid+g
# breakeven marker
bex=bx0+5*(cwid+g)-0.06
box(s,bex,2.95,0.012,2.95,fill='focus',radius=0,shape=MSO_SHAPE.RECTANGLE)
be=box(s,bex-0.75,2.7,1.7,0.42,fill='focus',radius=0.4); shape_text(be,"Breakeven · M6",size=9,color='white',bold=True,align=PP_ALIGN.CENTER)
text(s,"₹ Lakh / month",1.1,2.85,2,0.3,size=9.5,color='mute',font=BODY,bold=True)
# right KPI stack
kpis=[("~₹4–5 Cr","Year-1 revenue enabled (illustrative)"),("Month 6","Contribution breakeven"),
      ("3–4×","Target blended ROAS at scale"),("> 30%","Repeat-purchase rate goal")]
ky=2.7
for v,l in kpis:
    box(s,9.15,ky,3.35,0.9,fill='paper2',line='line',line_w=1,radius=0.08)
    text(s,v,9.4,ky+0.13,3,0.45,size=21,color='ink',font=TITLE,bold=True,track=-20)
    text(s,l,9.4,ky+0.55,3,0.3,size=9,color='mute',font=BODY)
    ky+=1.0
footer(s,18)

# ============================================================= 19 KPIs
s=slide('paper2')
header(s,"Measurement","Reported monthly — against sales, not vanity.")
gv=[("CAC","Cost to acquire a customer, by channel & creative."),
    ("ROAS","Blended & channel return on ad spend."),
    ("Repeat %","Repeat-purchase & subscription rate — the D2C flywheel."),
    ("Run-rate","Monthly revenue & MRR vs the milestone ladder."),
    ("AOV / LTV","Basket size and lifetime value vs CAC."),
    ("Reach & saves","Content performance feeding the funnel."),
    ("Reviews","Rating volume & sentiment — trust that compounds."),
    ("Contribution","Unit economics & path to profit.")]
cw=2.86; gap=0.22; per=4
for i,(m,d) in enumerate(gv):
    col=i%per; row=i//per
    x=0.85+col*(cw+gap); y=2.7+row*1.75
    box(s,x,y,cw,1.55,fill='paper',line='line',line_w=1,radius=0.08,shadow=True)
    text(s,m,x+0.28,y+0.24,cw-0.5,0.5,size=19,color='volt',font=TITLE,bold=True)
    text(s,d,x+0.28,y+0.78,cw-0.56,0.7,size=10,color='mute',font=BODY,ls=1.2)
footer(s,19)

# ============================================================= 20 NEXT STEPS + CTA
s=slide('paper')
header(s,"Next Steps","From plan to launch.")
st=[("Lock scope & budget","Confirm the deliverables proposal, marketing budget and the go-to-market sequence."),
    ("Build the brand & store","Identity, packaging, Shopify store and the launch content engine."),
    ("Launch & reinvest","Go live on D2C, prove the funnel, and climb the milestone ladder.")]
cw=3.86; gap=0.22
for i,(h,b) in enumerate(st):
    x=0.85+i*(cw+gap)
    box(s,x,2.6,cw,1.9,fill='paper2',line='line',line_w=1,radius=0.08,shadow=True)
    text(s,f"0{i+1}",x+0.3,2.78,1,0.6,size=30,color='voltL',font=TITLE,bold=True)
    text(s,h,x+0.3,3.4,cw-0.6,0.4,size=15,color='ink',font=TITLE,bold=True)
    text(s,b,x+0.3,3.82,cw-0.6,0.6,size=10.5,color='mute',font=BODY,ls=1.2)
cta=box(s,0.85,4.85,11.63,1.75,fill='ink',radius=0.08,shadow=True)
text(s,"Let's give Mach Five its launch velocity.",1.25,5.15,9,0.9,size=26,color='white',font=TITLE,bold=True,track=-30)
text(s,"business@sagemedia.in",1.25,6.05,6,0.4,size=14,color='voltL',font=TITLE,bold=True)
text(s,"Sage Media · Jaipur · Mumbai · Dubai",8.2,6.1,4.3,0.3,size=10.5,color='muteL',font=BODY,align=PP_ALIGN.RIGHT,track=20)
footer(s,20)

# ============================================================= 21 APPENDIX — AI PROMPTS
s=slide('ink')
box(s,0,0,SW,Inches(0.12),fill='volt',radius=0,shape=MSO_SHAPE.RECTANGLE)
text(s,"APPENDIX",0.85,0.55,9,0.35,size=10.5,color='voltL',font=BODY,bold=True,track=200)
text(s,"AI image prompts — mood board.",0.85,0.95,11,0.8,size=26,color='white',font=TITLE,bold=True,track=-30)
text(s,"Copy-paste into Midjourney / Firefly / DALL·E to generate the real photography for the board above.",
     0.85,1.7,11,0.4,size=11,color='muteL',font=BODY)
prompts=[("PRODUCT HERO","Matte-black Mach Five creatine + vitamin-C drink can on seamless studio background, volt-green rim light, condensation, dramatic hero product photography, high contrast, minimal, editorial, 4k."),
         ("STRENGTH / GYM","Young Indian male athlete mid-explosive lift, sweat and chalk, moody gym lighting with volt-green accent, motion energy, cinematic, 85mm."),
         ("CITRUS / VIT-C","Extreme macro of fresh citrus and effervescent bubbles, water splash, vibrant, clean studio, amber + green palette, science-forward."),
         ("BRAIN / FOCUS","Abstract glowing neural network, green-blue synapse lines over deep black, macro, futuristic, clean, cognition energy."),
         ("RECOVERY / CALM","Post-workout Indian athlete resting with towel and drink, soft natural window light, warm authentic lifestyle, calm."),
         ("LIFESTYLE","Diverse young Indians laughing after a run, urban rooftop at golden hour, candid, energetic, aspirational, cinematic.")]
cw=3.86; gap=0.22; per=3
for i,(h,p) in enumerate(prompts):
    col=i%per; row=i//per
    x=0.85+col*(cw+gap); y=2.25+row*2.15
    box(s,x,y,cw,1.95,fill='ink2',radius=0.06)
    box(s,x,y,cw,0.07,fill='volt',radius=0,shape=MSO_SHAPE.RECTANGLE)
    text(s,h,x+0.26,y+0.22,cw-0.5,0.3,size=10.5,color='voltL',font=BODY,bold=True,track=60)
    text(s,'“'+p+'”',x+0.26,y+0.6,cw-0.5,1.2,size=9.5,color='muteL',font=BODY,ls=1.25)
footer(s,21)

out='/home/user/Sage-portfolio/MachFive-GTM-Growth-Plan-SageMedia.pptx'
prs.save(out); print("saved",out,"slides",len(prs.slides._sldIdLst))
