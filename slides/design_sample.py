#!/usr/bin/env python3
"""Sage Media — proposal deck DESIGN DIRECTION sample (placeholder content).
Elevated creative-pitch-deck look on our format + fonts (Glacial Indifference
+ Inter). 16:9. Native, editable shapes only. Accent/imagery swap per brand."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

C = {
    'paper':'F4F1EA','paper2':'ECE7DB','paper3':'E3DCCB',
    'ink':'16171C','ink2':'22242C','ink3':'2E313B',
    'accent':'FF5A30','accent2':'E9A94A','accdim':'6E2A18',
    'mute':'7C7566','muteL':'B7AF9E','line':'D9D1C0','white':'FFFFFF',
}
def rgb(h): return RGBColor.from_string(C.get(h,h))
TITLE='Glacial Indifference'; BODY='Inter'

prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
BLANK=prs.slide_layouts[6]; SW,SH=prs.slide_width,prs.slide_height

def slide(bg='paper'):
    s=prs.slides.add_slide(BLANK)
    r=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,0,0,SW,SH)
    r.fill.solid(); r.fill.fore_color.rgb=rgb(bg); r.line.fill.background(); r.shadow.inherit=False
    return s

def _font(run,size,color,font,bold,italic=False,spc=None):
    run.font.size=Pt(size); run.font.bold=bold; run.font.italic=italic
    run.font.name=font; run.font.color.rgb=rgb(color)
    rPr=run._r.get_or_add_rPr()
    if spc is not None: rPr.set('spc',str(int(spc)))
    for tag in ('a:latin','a:cs'):
        e=rPr.find(qn(tag))
        if e is None: e=rPr.makeelement(qn(tag),{}); rPr.append(e)
        e.set('typeface',font)

def text(s,txt,l,t,w,h,size=14,color='ink',font=BODY,bold=False,align=PP_ALIGN.LEFT,
         anchor=MSO_ANCHOR.TOP,italic=False,ls=1.05,track=None,sp_after=2):
    tb=s.shapes.add_textbox(Inches(l),Inches(t),Inches(w),Inches(h))
    tf=tb.text_frame; tf.word_wrap=True; tf.vertical_anchor=anchor
    tf.margin_left=0;tf.margin_right=0;tf.margin_top=0;tf.margin_bottom=0
    for i,ln in enumerate(txt.split('\n')):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.alignment=align; p.space_after=Pt(sp_after); p.space_before=Pt(0)
        try: p.line_spacing=ls
        except Exception: pass
        r=p.add_run(); r.text=ln; _font(r,size,color,font,bold,italic,spc=track)
    return tb

def box(s,l,t,w,h,fill='white',line=None,line_w=1.0,radius=0.08,shape=MSO_SHAPE.ROUNDED_RECTANGLE,shadow=False):
    sp=s.shapes.add_shape(shape,Inches(l),Inches(t),Inches(w),Inches(h))
    if fill is None: sp.fill.background()
    else: sp.fill.solid(); sp.fill.fore_color.rgb=rgb(fill)
    if line is None: sp.line.fill.background()
    else: sp.line.color.rgb=rgb(line); sp.line.width=Pt(line_w)
    sp.shadow.inherit=False
    if shape==MSO_SHAPE.ROUNDED_RECTANGLE:
        try: sp.adjustments[0]=radius
        except Exception: pass
    if shadow:
        el=sp._element.spPr; ef=el.makeelement(qn('a:effectLst'),{})
        sh=el.makeelement(qn('a:outerShdw'),{'blurRad':'140000','dist':'50000','dir':'5400000','rotWithShape':'0'})
        clr=el.makeelement(qn('a:srgbClr'),{'val':'16171C'}); al=el.makeelement(qn('a:alpha'),{'val':'20000'})
        clr.append(al); sh.append(clr); ef.append(sh); el.append(ef)
    return sp

def grad(sp,c1,c2,angle=55):
    """duotone linear gradient fill for image-placeholder blocks"""
    spPr=sp._element.spPr
    for tag in ('a:solidFill','a:noFill','a:gradFill','a:blipFill','a:pattFill'):
        e=spPr.find(qn(tag))
        if e is not None: spPr.remove(e)
    g=spPr.makeelement(qn('a:gradFill'),{}); lst=spPr.makeelement(qn('a:gsLst'),{})
    for pos,c in [(0,c1),(100000,c2)]:
        gs=spPr.makeelement(qn('a:gs'),{'pos':str(pos)}); clr=spPr.makeelement(qn('a:srgbClr'),{'val':C.get(c,c)})
        gs.append(clr); lst.append(gs)
    g.append(lst); g.append(spPr.makeelement(qn('a:lin'),{'ang':str(int(angle*60000)),'scaled':'1'}))
    ln=spPr.find(qn('a:ln'))
    if ln is not None: spPr.insert(list(spPr).index(ln),g)
    else: spPr.append(g)
    return sp

def shape_text(sp,txt,size=14,color='white',font=BODY,bold=False,align=PP_ALIGN.LEFT,anchor=MSO_ANCHOR.MIDDLE,track=None):
    tf=sp.text_frame; tf.word_wrap=True; tf.vertical_anchor=anchor
    tf.margin_left=Inches(0.14);tf.margin_right=Inches(0.14);tf.margin_top=Inches(0.04);tf.margin_bottom=Inches(0.04)
    for i,ln in enumerate(txt.split('\n')):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph(); p.alignment=align
        r=p.add_run(); r.text=ln; _font(r,size,color,font,bold,spc=track)
    return sp

def imgblock(s,l,t,w,h,c1='ink2',c2='accent',label='BRAND IMAGERY',angle=55,radius=0.06):
    b=box(s,l,t,w,h,fill='ink2',radius=radius); grad(b,c1,c2,angle)
    # camera/frame glyph centered (simple geometric): a ring + dot
    cx=l+w/2; cy=t+h/2
    box(s,cx-0.36,cy-0.36,0.72,0.72,fill=None,line='white',line_w=1.6,radius=0.5,shape=MSO_SHAPE.OVAL).fill.background()
    box(s,cx-0.09,cy-0.09,0.18,0.18,fill='white',radius=0.5,shape=MSO_SHAPE.OVAL)
    # label chip bottom-left
    chip=box(s,l+0.28,t+h-0.62,1.9,0.36,fill='ink',radius=0.5)
    chip.fill.fore_color.rgb=rgb('16171C')
    shape_text(chip,label,size=8.5,color='white',bold=True,align=PP_ALIGN.CENTER,track=60)
    return b

def eyebrow(s,txt,l=0.85,t=0.72,dark=False,color=None):
    col=color or ('accent2' if dark else 'accent')
    tb=text(s,txt.upper(),l,t,7,0.35,size=10.5,color=col,font=BODY,bold=True,track=200)
    return tb

def footer(s,idx,total=7,dark=False):
    col='muteL' if dark else 'mute'
    box(s,0.85,7.02,11.63,0.014,fill=('ink3' if dark else 'line'),radius=0,shape=MSO_SHAPE.RECTANGLE)
    text(s,"SAGE MEDIA",0.85,7.1,4,0.3,size=8.5,color=col,font=BODY,bold=True,track=120)
    text(s,f"{idx:02d} / {total:02d}",11.0,7.1,1.48,0.3,size=8.5,color=col,font=BODY,align=PP_ALIGN.RIGHT,track=80)

def sectnum(s,n,color='ink2'):
    """oversized ghost index number as editorial motif"""
    text(s,n,8.1,-0.6,6,6,size=300,color=color,font=TITLE,bold=True,align=PP_ALIGN.LEFT,track=-60)

# ============ 1 · COVER ============
s=slide('paper')
box(s,0,0,0.22,SH,fill='accent',radius=0,shape=MSO_SHAPE.RECTANGLE)   # left spine
# right imagery (two stacked editorial blocks)
imgblock(s,8.55,0.0,4.78,4.55,'ink','accent',angle=60,radius=0.0)
imgblock(s,8.55,4.55,4.78,2.95,'accent','ink2',label='LIFESTYLE',angle=35,radius=0.0)
eyebrow(s,"Business Proposal · 2026",l=0.95,t=0.85)
text(s,"Aurora",0.9,1.5,8,1.5,size=76,color='ink',font=TITLE,bold=True,track=-60)
box(s,0.98,3.05,0.9,0.09,fill='accent',radius=0,shape=MSO_SHAPE.RECTANGLE)
text(s,"A brand growth & marketing proposal —\nstrategy, creative and a measurable plan of action.",
     0.95,3.4,6.9,1.1,size=15.5,color='ink3',font=BODY,ls=1.3)
text(s,"Prepared for  ·  [ Client Name ]",0.95,5.7,6,0.35,size=11,color='mute',font=BODY,bold=True,track=40)
text(s,"Prepared by Sage Media",0.95,6.55,6,0.35,size=12.5,color='ink',font=TITLE,bold=True)
text(s,"Jaipur · Mumbai · Dubai  ·  business@sagemedia.in",0.95,6.95,7,0.3,size=9.5,color='mute',font=BODY,track=30)

# ============ 2 · SECTION DIVIDER (dark) ============
s=slide('ink')
sectnum(s,"01",color='ink3')
box(s,0.85,0,0.22,SH,fill='accent',radius=0,shape=MSO_SHAPE.RECTANGLE)
text(s,"SECTION 01",1.35,2.35,6,0.4,size=11,color='accent2',font=BODY,bold=True,track=220)
text(s,"The\nOpportunity",1.3,2.85,8,2.4,size=62,color='white',font=TITLE,bold=True,track=-40,ls=0.98)
text(s,"Where the brand is today, the gap in the market, and the moment to move.",
     1.35,5.35,7.4,0.7,size=14.5,color='muteL',font=BODY,ls=1.3)
footer(s,2,dark=True)

# ============ 3 · OVERVIEW / STATS ============
s=slide('paper')
eyebrow(s,"At a glance")
text(s,"A brand with momentum —\nand room to run.",0.85,1.15,8.5,1.5,size=34,color='ink',font=TITLE,bold=True,track=-30,ls=1.0)
text(s,"A short, editorial framing paragraph goes here — two lines that set up the numbers below and give the reader the story before the data.",
     0.85,2.75,7.4,0.9,size=13.5,color='ink3',font=BODY,ls=1.35)
imgblock(s,9.9,1.1,2.58,3.0,'ink2','accent',label='HERO SHOT',angle=50)
# stat row
stats=[("3.2×","projected ROI"),("+48%","audience growth"),("12","markets"),("₹—","your number")]
sx=0.85; sw=2.85; base=4.55
for i,(v,l) in enumerate(stats):
    x=sx+i*sw
    if i>0: box(s,x-0.18,base+0.1,0.014,1.15,fill='line',radius=0,shape=MSO_SHAPE.RECTANGLE)
    text(s,v,x,base,sw-0.3,0.8,size=46,color='accent',font=TITLE,bold=True,track=-30)
    text(s,l.upper(),x,base+0.95,sw-0.3,0.4,size=10,color='mute',font=BODY,bold=True,track=80)
text(s,"Every figure is illustrative — swapped for the brand's real data at kickoff.",
     0.85,6.35,9,0.3,size=9.5,color='mute',font=BODY,italic=True)
footer(s,3)

# ============ 4 · APPROACH (two column) ============
s=slide('paper')
eyebrow(s,"Our approach")
text(s,"Three moves, one\ndirection of travel.",0.85,1.15,7,1.5,size=34,color='ink',font=TITLE,bold=True,track=-30,ls=1.0)
pts=[("01","Understand","Audit the brand, audience and category — find the sharpest angle to own."),
     ("02","Create","Build a distinctive system: identity, messaging and a content engine that scales."),
     ("03","Grow","Launch, measure and optimise against outcomes — not vanity metrics.")]
py=3.05
for n,h,b in pts:
    text(s,n,0.85,py,0.9,0.6,size=26,color='accent',font=TITLE,bold=True)
    text(s,h,1.75,py+0.02,5,0.4,size=16,color='ink',font=TITLE,bold=True)
    text(s,b,1.75,py+0.46,5.0,0.6,size=11.5,color='ink3',font=BODY,ls=1.25)
    if n!="03": box(s,1.75,py+1.12,5.0,0.012,fill='line',radius=0,shape=MSO_SHAPE.RECTANGLE)
    py+=1.28
imgblock(s,8.5,1.15,4.0,5.35,'accent','ink','APPROACH',angle=60)
footer(s,4)

# ============ 5 · DATA ============
s=slide('paper2')
eyebrow(s,"The opportunity in numbers")
text(s,"The case, in one chart.",0.85,1.15,9,0.9,size=34,color='ink',font=TITLE,bold=True,track=-30)
# bar chart card
box(s,0.85,2.35,7.0,4.1,fill='paper',line='line',line_w=1,radius=0.05,shadow=True)
text(s,"Channel contribution (illustrative)",1.2,2.6,5.5,0.4,size=12.5,color='ink',font=TITLE,bold=True)
bars=[("Social",92,'accent'),("Search",70,'ink2'),("Content",58,'ink2'),("Influencer",44,'accent'),("Email",30,'ink2')]
by=3.3; bh=0.34; bgap=0.2; bx=2.5; bfull=4.6
for nm,v,col in bars:
    text(s,nm,1.2,by-0.02,1.2,0.35,size=10,color='ink3',font=BODY,bold=True,align=PP_ALIGN.RIGHT,anchor=MSO_ANCHOR.MIDDLE)
    box(s,bx,by,bfull,bh,fill='paper3',radius=0.4)
    box(s,bx,by,max(bfull*v/92,0.1),bh,fill=col,radius=0.4)
    text(s,str(v),bx+bfull*v/92+0.1,by-0.02,0.6,0.35,size=10,color='ink',font=TITLE,bold=True,anchor=MSO_ANCHOR.MIDDLE)
    by+=bh+bgap
# highlight stat card (right)
hc=box(s,8.1,2.35,4.4,4.1,fill='ink',radius=0.05,shadow=True)
text(s,"HEADLINE METRIC",8.45,2.7,3.7,0.35,size=10,color='accent2',font=BODY,bold=True,track=180)
text(s,"3.2×",8.45,3.15,3.7,1.4,size=78,color='white',font=TITLE,bold=True,track=-40)
text(s,"return on marketing spend,\nmodelled conservatively.",8.45,4.75,3.7,0.7,size=13,color='muteL',font=BODY,ls=1.3)
box(s,8.45,5.75,3.6,0.5,fill='ink3',radius=0.4)
box(s,8.45,5.75,3.6*0.62,0.5,fill='accent',radius=0.4)
text(s,"62% to target",8.55,5.79,3,0.42,size=10,color='white',font=BODY,bold=True,anchor=MSO_ANCHOR.MIDDLE)
footer(s,5)

# ============ 6 · PLAN / PHASES ============
s=slide('paper')
eyebrow(s,"The plan")
text(s,"A phased 90-day rollout.",0.85,1.15,9,0.9,size=34,color='ink',font=TITLE,bold=True,track=-30)
phases=[("Phase 01","Foundation","Brand system, messaging &\nchannel setup.",'accent'),
        ("Phase 02","Momentum","Always-on content, campaigns\n& community building.",'ink2'),
        ("Phase 03","Scale","Performance, partnerships\n& optimisation to target.",'accent')]
cw=3.78; gap=0.24; x0=0.85; y0=2.5; ch=3.8
for i,(ph,h,b,acc) in enumerate(phases):
    x=x0+i*(cw+gap)
    dark = (acc=='ink2')
    box(s,x,y0,cw,ch,fill=('ink' if dark else 'paper'),line=(None if dark else 'line'),line_w=1,radius=0.05,shadow=True)
    text(s,ph.upper(),x+0.35,y0+0.4,cw-0.7,0.35,size=10.5,color=('accent2' if dark else 'accent'),font=BODY,bold=True,track=140)
    text(s,("0"+str(i+1)),x+0.32,y0+0.75,cw-0.7,1.0,size=56,color=('ink3' if dark else 'paper3'),font=TITLE,bold=True)
    text(s,h,x+0.35,y0+2.0,cw-0.7,0.5,size=20,color=('white' if dark else 'ink'),font=TITLE,bold=True)
    text(s,b,x+0.35,y0+2.6,cw-0.7,0.9,size=11.5,color=('muteL' if dark else 'ink3'),font=BODY,ls=1.3)
    box(s,x,y0+ch-0.09,cw,0.09,fill=('accent2' if dark else 'accent'),radius=0,shape=MSO_SHAPE.RECTANGLE)
footer(s,6)

# ============ 7 · CLOSING (dark) ============
s=slide('ink')
box(s,0.85,0,0.22,SH,fill='accent',radius=0,shape=MSO_SHAPE.RECTANGLE)
# small stacked motif
for i,(yy,cc) in enumerate([(0.95,'ink3'),(1.16,'ink3'),(1.37,'accent'),(1.58,'accent2')]):
    box(s,11.0,yy,1.5,0.13,fill=cc,radius=0.3)
text(s,"Let's build a brand\nworth remembering.",1.35,2.55,10.6,2.2,size=52,color='white',font=TITLE,bold=True,track=-40,ls=1.02)
box(s,1.4,4.75,0.9,0.09,fill='accent',radius=0,shape=MSO_SHAPE.RECTANGLE)
text(s,"business@sagemedia.in",1.35,5.35,7,0.5,size=17,color='white',font=TITLE,bold=True)
text(s,"Sage Media · Jaipur · Mumbai · Dubai",1.35,5.95,8,0.4,size=11.5,color='muteL',font=BODY,track=30)
text(s,"SAMPLE — DESIGN DIRECTION ONLY · PLACEHOLDER CONTENT",1.35,6.75,10,0.3,size=8.5,color='ink3',font=BODY,bold=True,track=120)

out='/home/user/Sage-portfolio/Sage-Proposal-Design-Sample.pptx'
prs.save(out); print("saved",out)
